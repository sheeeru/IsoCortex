"""
IsoCortex — ingestion/extractor.py
====================================
Format-specific text extraction engine.

Responsibilities (FR-1, FR-1b):
  - Accept a ScannedFile from scanner.py and dispatch to the correct
    format-specific extractor based on format_category and extension.
  - Extract clean, normalised plain-text from every supported format.
  - Apply FR-1b structured data linearization rules for spreadsheets,
    presentations, word documents, email, and HTML.
  - Return structured ExtractionResult objects — never raw strings.
  - Handle every known failure mode per format gracefully.
  - Emit structured log messages at every decision point.

Supported formats (Section 1.3):
  plain_text   -> .txt  .md  .log
  source_code  -> .py   .cpp .c   .h   .js  .ts  .java
  pdf          -> .pdf
  word         -> .docx .odt
  presentation -> .pptx .odp
  spreadsheet  -> .xlsx .xls .ods .csv
  email        -> .eml
  web          -> .html .htm

SRS References: FR-1, FR-1b, FR-6, NFR-2, NFR-3, NFR-4,
                ASM-3, ASM-4, CON-3, OOS-5

Author : Shaheer Qureshi
Project: IsoCortex
"""

from __future__ import annotations

import csv
import email as email_lib
import email.policy as _email_policy
import hashlib
import html as html_lib
import io
import logging
import re
import time
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Optional

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
MAX_CHARS_PER_DOC: int          = 5_000_000   # ~5 MB of text per chunk
MIN_CONTENT_LENGTH: int         = 3
TEXT_ENCODINGS: tuple[str, ...] = (
    "utf-8",
    "utf-8-sig",
    "latin-1",
    "cp1252",
    "iso-8859-1",
)
HEADER_SCAN_ROWS: int           = 5           # Rows to scan for best header row


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class ExtractedChunk:
    """
    A single logical unit of text extracted from one document or
    sub-document.

    Attributes
    ----------
    text         : str — Normalised plain-text content.
    source_label : str — Human-readable origin label
                         e.g. "Sheet: Revenue", "Slide 3", "Page 1".
    char_count   : int — Number of characters in text.
    word_count   : int — Approximate word count.
    """
    text:         str
    source_label: str
    char_count:   int = field(init=False)
    word_count:   int = field(init=False)

    def __post_init__(self) -> None:
        self.char_count = len(self.text)
        self.word_count = len(self.text.split())


@dataclass
class ExtractionResult:
    """
    Top-level return value of :func:`extract`.

    Attributes
    ----------
    absolute_path   : Path                 — Source file path.
    format_category : str                  — e.g. 'pdf', 'spreadsheet'.
    extension       : str                  — e.g. '.xlsx'.
    chunks          : list[ExtractedChunk] — Extracted content units.
    success         : bool                 — False if extraction failed.
    error_message   : str | None           — Populated when success=False.
    elapsed_seconds : float                — Wall-clock extraction time.
    file_hash_md5   : str                  — MD5 of raw file bytes.
    """
    absolute_path:   Path
    format_category: str
    extension:       str
    chunks:          list[ExtractedChunk] = field(default_factory=list)
    success:         bool                 = True
    error_message:   Optional[str]        = None
    elapsed_seconds: float                = 0.0
    file_hash_md5:   str                  = ""

    @property
    def full_text(self) -> str:
        return "\n\n".join(c.text for c in self.chunks)

    @property
    def total_words(self) -> int:
        return sum(c.word_count for c in self.chunks)

    @property
    def total_chars(self) -> int:
        return sum(c.char_count for c in self.chunks)

    @property
    def chunk_count(self) -> int:
        return len(self.chunks)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract(scanned_file: object) -> ExtractionResult:
    """
    Dispatch *scanned_file* to the correct format extractor.

    Extension-level overrides are checked first so that non-OOXML
    formats (.odt, .odp, .xls, .ods) route to the correct library
    instead of being sent to openpyxl / python-docx / python-pptx.

    Parameters
    ----------
    scanned_file : ScannedFile
        A ScannedFile instance produced by scan_directory().

    Returns
    -------
    ExtractionResult
        Always returns a result object. On failure success=False and
        error_message is populated. This function never raises.
    """
    path:     Path = scanned_file.absolute_path      # type: ignore[attr-defined]
    category: str  = scanned_file.format_category    # type: ignore[attr-defined]
    ext:      str  = scanned_file.extension          # type: ignore[attr-defined]

    logger.debug(
        "[EXTRACTOR] Starting  file=%s  category=%s  ext=%s",
        path.name, category, ext,
    )

    result = ExtractionResult(
        absolute_path   = path,
        format_category = category,
        extension       = ext,
    )

    t_start = time.perf_counter()

    try:
        result.file_hash_md5 = _md5_of_file(path)

        # ── Extension-level overrides (non-OOXML formats) ─────────────
        ext_lower = ext.lower()
        if ext_lower == ".odt":
            extractor_fn = _extract_odt
        elif ext_lower == ".odp":
            extractor_fn = _extract_odp
        elif ext_lower == ".xls":
            extractor_fn = _extract_xls
        elif ext_lower == ".ods":
            extractor_fn = _extract_ods
        else:
            extractor_fn = _DISPATCH.get(category)

        if extractor_fn is None:
            _fail(result, f"No extractor registered for category '{category}'")
        else:
            extractor_fn(path, result)

    except MemoryError:
        _fail(result, "MemoryError — file too large to process")
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"{type(exc).__name__}: {exc}")

    result.elapsed_seconds = time.perf_counter() - t_start

    if result.success and result.chunk_count == 0:
        _fail(result, "Extraction produced zero chunks")

    if result.success:
        logger.info(
            "[EXTRACTOR] OK      %-40s  chunks=%d  words=%d  elapsed=%.3fs",
            path.name, result.chunk_count,
            result.total_words, result.elapsed_seconds,
        )
    else:
        logger.warning(
            "[EXTRACTOR] FAILED  %-40s  reason=%s",
            path.name, result.error_message,
        )

    return result


def extract_batch(scanned_files: list) -> list[ExtractionResult]:
    """
    Extract text from a list of ScannedFile objects.

    Failures on individual files do not stop the batch.
    Progress is logged at DEBUG per-file and INFO every 100 files
    to avoid flooding stdout on large directories.

    Parameters
    ----------
    scanned_files : list[ScannedFile]

    Returns
    -------
    list[ExtractionResult]
        Same length as input.
    """
    results: list[ExtractionResult] = []
    total = len(scanned_files)

    for idx, sf in enumerate(scanned_files, start=1):
        if idx % 100 == 0 or idx == total:
            logger.info("[EXTRACTOR] Batch progress: %d/%d", idx, total)
        logger.debug(
            "[EXTRACTOR] Batch %d/%d  file=%s",
            idx, total,
            sf.absolute_path.name,               # type: ignore[attr-defined]
        )
        results.append(extract(sf))

    succeeded = sum(1 for r in results if r.success)
    failed    = total - succeeded
    logger.info(
        "[EXTRACTOR] Batch complete  total=%d  succeeded=%d  failed=%d",
        total, succeeded, failed,
    )
    return results


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _fail(result: ExtractionResult, message: str) -> None:
    """Mark result as failed, clear chunks, and log a warning."""
    result.success       = False
    result.error_message = message
    result.chunks        = []


def _add_chunk(
    result:              ExtractionResult,
    raw_text:            str,
    source_label:        str,
    preserve_formatting: bool = False,
) -> None:
    """
    Normalise raw_text, enforce MAX_CHARS_PER_DOC, append to result.
    Discards chunks shorter than MIN_CONTENT_LENGTH.

    Parameters
    ----------
    preserve_formatting : bool
        If True, leading whitespace (indentation) is preserved.
        Use for source code files where indentation is meaningful.
    """
    if preserve_formatting:
        text = _normalise_text_preserve_indent(raw_text)
    else:
        text = _normalise_text(raw_text)

    if len(text) > MAX_CHARS_PER_DOC:
        logger.warning(
            "[EXTRACTOR] Chunk truncated at %d chars (was %d)  label=%s",
            MAX_CHARS_PER_DOC, len(text), source_label,
        )
        text = text[:MAX_CHARS_PER_DOC]

    if len(text) < MIN_CONTENT_LENGTH:
        logger.debug(
            "[EXTRACTOR] Chunk too short (%d chars), discarded  label=%s",
            len(text), source_label,
        )
        return

    result.chunks.append(ExtractedChunk(text=text, source_label=source_label))


def _normalise_text(text: str) -> str:
    """
    Consistent text normalisation applied across all format extractors.

    Steps:
      1. Unicode NFC normalisation.
      2. Replace non-breaking spaces and zero-width characters.
      3. Normalise all line endings to LF.
      4. Strip leading/trailing whitespace per line.
      5. Collapse 3+ consecutive blank lines to exactly 2.
      6. Strip overall leading/trailing whitespace.
    """
    if not text:
        return ""

    text = unicodedata.normalize("NFC", text)
    text = (
        text
        .replace("\u00a0", " ")
        .replace("\u200b", "")
        .replace("\u200c", "")
        .replace("\u200d", "")
        .replace("\ufeff", "")
        .replace("\r\n",  "\n")
        .replace("\r",    "\n")
    )

    lines = [line.strip() for line in text.split("\n")]

    collapsed:   list[str] = []
    blank_count: int        = 0
    for line in lines:
        if line == "":
            blank_count += 1
            if blank_count <= 2:
                collapsed.append(line)
        else:
            blank_count = 0
            collapsed.append(line)

    return "\n".join(collapsed).strip()


def _normalise_text_preserve_indent(text: str) -> str:
    """
    Normalise text while preserving leading whitespace (indentation).

    Used for source code files where indentation is semantically
    meaningful (e.g. Python blocks, C++ brace style).

    Steps:
      1. Unicode NFC normalisation.
      2. Replace non-breaking spaces and zero-width characters.
      3. Normalise all line endings to LF.
      4. Strip *trailing* whitespace per line (keep leading).
      5. Collapse 3+ consecutive blank lines to exactly 2.
      6. Strip overall leading/trailing whitespace.
    """
    if not text:
        return ""

    text = unicodedata.normalize("NFC", text)
    text = (
        text
        .replace("\u00a0", " ")
        .replace("\u200b", "")
        .replace("\u200c", "")
        .replace("\u200d", "")
        .replace("\ufeff", "")
        .replace("\r\n",  "\n")
        .replace("\r",    "\n")
    )

    # Strip trailing whitespace only — preserve leading indentation
    lines = [line.rstrip() for line in text.split("\n")]

    collapsed:   list[str] = []
    blank_count: int        = 0
    for line in lines:
        if line == "":
            blank_count += 1
            if blank_count <= 2:
                collapsed.append(line)
        else:
            blank_count = 0
            collapsed.append(line)

    return "\n".join(collapsed).strip()


def _md5_of_file(path: Path, chunk_size: int = 65_536) -> str:
    """Compute MD5 hash of a file in streaming chunks (memory-safe)."""
    hasher = hashlib.md5()
    try:
        with open(path, "rb") as fh:
            while True:
                block = fh.read(chunk_size)
                if not block:
                    break
                hasher.update(block)
        return hasher.hexdigest()
    except OSError as exc:
        logger.debug(
            "[EXTRACTOR] MD5 hash failed for %s: %s", path.name, exc,
        )
        return ""


def _read_text_with_fallback(path: Path) -> str:
    """
    Read a text file trying multiple encodings (ASM-3).
    Never raises — uses replace error handler as last resort.
    """
    for enc in TEXT_ENCODINGS:
        try:
            return path.read_text(encoding=enc, errors="strict")
        except (UnicodeDecodeError, LookupError):
            continue
    logger.warning(
        "[EXTRACTOR] Encoding fallback (utf-8 replace) for %s", path.name,
    )
    return path.read_text(encoding="utf-8", errors="replace")


def _safe_get_email_content(part: object) -> str:
    """
    Safely extract string content from an email message part.

    email.policy.default makes get_content() return str | bytes | list | dict
    depending on MIME type. This helper always returns a plain str.
    """
    try:
        raw = part.get_content()                     # type: ignore[attr-defined]
        if isinstance(raw, str):
            return raw
        if isinstance(raw, bytes):
            return raw.decode("utf-8", errors="replace")
        return str(raw)
    except Exception:                                # pylint: disable=broad-except
        return ""


# ---------------------------------------------------------------------------
# Format extractors — Plain Text & Source Code
# ---------------------------------------------------------------------------

def _extract_plain_text(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .txt / .md / .log / source code files.

    Tries multiple encodings in order (ASM-3).
    Uses 'replace' error handler as final fallback.
    Preserves indentation for source code files.
    """
    text = _read_text_with_fallback(path)
    is_source = result.format_category == "source_code"
    logger.debug(
        "[EXTRACTOR] PlainText  file=%s  chars=%d  preserve_indent=%s",
        path.name, len(text), is_source,
    )
    _add_chunk(
        result,
        text,
        source_label=f"File: {path.name}",
        preserve_formatting=is_source,
    )


# ---------------------------------------------------------------------------
# Format extractors — PDF
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from PDF using PyMuPDF (fitz).

    Each page produces one ExtractedChunk labelled "Page N" for
    fine-grained traceability in search results (FR-6).

    Handles:
      - Password-protected PDFs.
      - Scanned image PDFs with no text (OOS-5).
      - Corrupted per-page streams.
      - PDFs that report is_encrypted=False but are actually encrypted.
    """
    try:
        import fitz                                  # type: ignore[import-untyped]
    except ImportError:
        _fail(result, "pymupdf not installed — run: pip install pymupdf")
        return

    try:
        doc = fitz.open(str(path))
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot open PDF: {exc}")
        return

    if doc.is_encrypted:
        doc.close()
        _fail(result, "PDF is password-protected — cannot extract text")
        return

    # Additional guard: some encrypted PDFs pass is_encrypted check
    try:
        _ = doc[0]
    except Exception:
        doc.close()
        _fail(
            result,
            "PDF appears to be encrypted or corrupted — cannot access pages",
        )
        return

    total_pages     = doc.page_count
    extracted_pages = 0

    for page_num in range(total_pages):
        try:
            page = doc[page_num]
            raw: str = page.get_text("text")         # type: ignore[arg-type]
            if raw and raw.strip():
                _add_chunk(
                    result,
                    raw,
                    source_label=f"Page {page_num + 1}",
                )
                extracted_pages += 1
        except Exception as exc:                     # pylint: disable=broad-except
            logger.warning(
                "[EXTRACTOR] PDF page %d failed  file=%s  reason=%s",
                page_num + 1, path.name, exc,
            )
            continue

    doc.close()

    if extracted_pages == 0:
        _fail(
            result,
            "PDF contains no extractable text. "
            "It may be a scanned image PDF — OCR is OOS-5 for v1.0.",
        )
        return

    logger.debug(
        "[EXTRACTOR] PDF  file=%s  pages=%d/%d extracted",
        path.name, extracted_pages, total_pages,
    )


# ---------------------------------------------------------------------------
# Format extractors — Word Documents
# ---------------------------------------------------------------------------

def _extract_docx(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .docx files using python-docx.

    FR-1b Word rule:
      - Body paragraphs in document order.
      - Heading paragraphs prefixed with style name.
      - Table cells row-by-row, pipe-separated.
      - Section headers/footers appended.

    Note: This extractor handles .docx (OOXML) only.
    For .odt files, _extract_odt is used instead.
    """
    try:
        from docx import Document                    # type: ignore[import-untyped]
    except ImportError:
        _fail(result, "python-docx not installed — run: pip install python-docx")
        return

    try:
        doc = Document(str(path))
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot open Word document: {exc}")
        return

    lines: list[str] = []

    # ---- Body paragraphs ------------------------------------------------
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name: str = ""
        if para.style is not None and para.style.name is not None:
            style_name = str(para.style.name)
        if "Heading" in style_name:
            lines.append(f"[{style_name}] {text}")
        else:
            lines.append(text)

    # ---- Tables (FR-1b Word) --------------------------------------------
    for table_idx, table in enumerate(doc.tables, start=1):
        lines.append(f"[Table {table_idx}]")
        for row in table.rows:
            row_cells = [
                cell.text.strip()
                for cell in row.cells
                if cell.text.strip()
            ]
            if row_cells:
                lines.append(" | ".join(row_cells))

    # ---- Headers & Footers ----------------------------------------------
    for section in doc.sections:
        try:
            for hf_para in (
                list(section.header.paragraphs) +
                list(section.footer.paragraphs)
            ):
                hf_text = hf_para.text.strip()
                if hf_text:
                    lines.append(f"[Header/Footer] {hf_text}")
        except Exception:                            # pylint: disable=broad-except
            pass

    if not lines:
        _fail(result, "Word document contains no extractable text")
        return

    _add_chunk(result, "\n".join(lines), source_label=f"Document: {path.name}")
    logger.debug("[EXTRACTOR] DOCX  file=%s  lines=%d", path.name, len(lines))


def _extract_odt(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .odt files using odfpy.

    FR-1b Word rule (ODT variant):
      - All <text:p> paragraphs extracted in document order.
      - All <text:h> headings extracted and prefixed with level.
      - Tables extracted row-by-row, pipe-separated.
      - List items extracted with bullet markers.

    Note: python-docx does NOT support .odt. odfpy is required.
    """
    try:
        from odf.opendocument import load as odf_load                     # type: ignore[import-untyped]
        from odf.text import P as OdfP, H as OdfH, ListItem as OdfListItem  # type: ignore[import-untyped]
        from odf.table import (                                           # type: ignore[import-untyped]
            Table as OdfTable,
            TableRow as OdfTableRow,
            TableCell as OdfTableCell,
        )
    except ImportError:
        _fail(result, "odfpy not installed — run: pip install odfpy")
        return

    try:
        doc = odf_load(str(path))
    except Exception as exc:
        _fail(result, f"Cannot open ODT document: {exc}")
        return

    lines: list[str] = []

    # ---- Headings -------------------------------------------------------
    for heading in doc.getElementsByType(OdfH):
        try:
            outline_level = heading.getAttribute("outlinelevel")
            level = int(str(outline_level)) if outline_level is not None else 1
        except (ValueError, TypeError):
            level = 1
        text = _odf_element_text(heading).strip()
        if text:
            lines.append(f"[Heading {level}] {text}")

    # ---- Paragraphs -----------------------------------------------------
    for para in doc.getElementsByType(OdfP):
        text = _odf_element_text(para).strip()
        if text:
            lines.append(text)

    # ---- List items -----------------------------------------------------
    for item in doc.getElementsByType(OdfListItem):
        text = _odf_element_text(item).strip()
        if text:
            lines.append(f"• {text}")

    # ---- Tables ---------------------------------------------------------
    table_idx = 0
    for table in doc.getElementsByType(OdfTable):
        table_idx += 1
        lines.append(f"[Table {table_idx}]")
        for row in table.getElementsByType(OdfTableRow):
            cells = []
            for cell in row.getElementsByType(OdfTableCell):
                cell_text = _odf_element_text(cell).strip()
                if cell_text:
                    cells.append(cell_text)
            if cells:
                lines.append(" | ".join(cells))

    if not lines:
        _fail(result, "ODT document contains no extractable text")
        return

    _add_chunk(result, "\n".join(lines), source_label=f"Document: {path.name}")
    logger.debug("[EXTRACTOR] ODT  file=%s  lines=%d", path.name, len(lines))


def _odf_element_text(element: object) -> str:
    """Recursively extract all text content from an ODF element."""
    parts: list[str] = []
    try:
        for child in element.childNodes:             # type: ignore[attr-defined]
            if hasattr(child, "data"):
                parts.append(str(child.data))        # type: ignore[attr-defined]
            elif hasattr(child, "childNodes"):
                parts.append(_odf_element_text(child))
    except Exception:
        pass
    return "".join(parts)


# ---------------------------------------------------------------------------
# Format extractors — Presentations
# ---------------------------------------------------------------------------

def _extract_pptx(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .pptx files using python-pptx.

    FR-1b Presentation rule:
      - Each slide → one ExtractedChunk labelled "Slide N".
      - Title placed first in each chunk.
      - All shapes extracted recursively (text frames, tables, groups).
      - Speaker notes appended under [Notes].
      - Empty slides skipped.

    Note: This extractor handles .pptx (OOXML) only.
    For .odp files, _extract_odp is used instead.
    """
    try:
        from pptx import Presentation               # type: ignore[import-untyped]
    except ImportError:
        _fail(result, "python-pptx not installed — run: pip install python-pptx")
        return

    try:
        prs = Presentation(str(path))
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot open presentation: {exc}")
        return

    total_slides     = len(prs.slides)
    extracted_slides = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        title_text              = ""

        try:
            if (
                slide.shapes.title is not None
                and slide.shapes.title.has_text_frame
            ):
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    slide_lines.append(f"[Title] {title_text}")
        except Exception:                            # pylint: disable=broad-except
            pass

        for shape in slide.shapes:
            _extract_pptx_shape_recursive(shape, slide_lines, title_text)

        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf is not None:
                    notes = notes_tf.text.strip()
                    if notes:
                        slide_lines.append(f"[Notes] {notes}")
        except Exception:                            # pylint: disable=broad-except
            pass

        if not slide_lines:
            logger.debug(
                "[EXTRACTOR] PPTX slide %d empty  file=%s",
                slide_num, path.name,
            )
            continue

        _add_chunk(
            result,
            "\n".join(slide_lines),
            source_label=f"Slide {slide_num}",
        )
        extracted_slides += 1

    if extracted_slides == 0:
        _fail(result, "Presentation contains no extractable text")
        return

    logger.debug(
        "[EXTRACTOR] PPTX  file=%s  slides=%d/%d extracted",
        path.name, extracted_slides, total_slides,
    )


def _extract_pptx_shape_recursive(
    shape:     object,
    lines:     list[str],
    skip_text: str = "",
) -> None:
    """
    Recursively extract text from a pptx shape.
    Handles TextFrame, Table, and GroupShape without MSO_SHAPE_TYPE enum.
    """
    try:
        if hasattr(shape, "shapes"):
            for child in shape.shapes:               # type: ignore[attr-defined]
                _extract_pptx_shape_recursive(child, lines, skip_text)
            return

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:  # type: ignore[attr-defined]
            for para in shape.text_frame.paragraphs: # type: ignore[attr-defined]
                text = para.text.strip()
                if text and text != skip_text:
                    lines.append(text)

        if hasattr(shape, "has_table") and shape.has_table:            # type: ignore[attr-defined]
            for row in shape.table.rows:             # type: ignore[attr-defined]
                cells = [
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text.strip()
                ]
                if cells:
                    lines.append(" | ".join(cells))

    except Exception:                                # pylint: disable=broad-except
        pass


def _extract_odp(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .odp presentation files using odfpy.

    FR-1b Presentation rule (ODP variant):
      - Each <draw:page> treated as a logical slide unit.
      - Text extracted from all <text:p> elements within each frame.
      - Frames prefixed with slide number for traceability.
      - Speaker notes extracted if present.
      - Fallback: if no draw:page structure, extract all paragraphs.

    Note: python-pptx does NOT support .odp. odfpy is required.
    """
    try:
        from odf.opendocument import load as odf_load                     # type: ignore[import-untyped]
        from odf.text import P as OdfP                                    # type: ignore[import-untyped]
        from odf.draw import (                                            # type: ignore[import-untyped]
            Page as OdfPage,
            Frame as OdfFrame,
            Notes as OdfNotes,
        )
    except ImportError:
        _fail(result, "odfpy not installed — run: pip install odfpy")
        return

    try:
        doc = odf_load(str(path))
    except Exception as exc:
        _fail(result, f"Cannot open ODP presentation: {exc}")
        return

    # Try to extract by draw:page (slides)
    slides = doc.getElementsByType(OdfPage)

    if slides:
        extracted_slides = 0
        for slide_num, page in enumerate(slides, start=1):
            slide_lines: list[str] = []

            # Extract text from all frames within this page
            frames = page.getElementsByType(OdfFrame)
            for frame in frames:
                frame_text = _odf_element_text(frame).strip()
                if frame_text:
                    slide_lines.append(frame_text)

            # Extract notes if present
            try:
                notes_elements = page.getElementsByType(OdfNotes)
                for notes_elem in notes_elements:
                    notes_text = _odf_element_text(notes_elem).strip()
                    if notes_text:
                        slide_lines.append(f"[Notes] {notes_text}")
            except Exception:
                pass

            if not slide_lines:
                continue

            _add_chunk(
                result,
                "\n".join(slide_lines),
                source_label=f"Slide {slide_num}",
            )
            extracted_slides += 1

        if extracted_slides > 0:
            logger.debug(
                "[EXTRACTOR] ODP  file=%s  slides=%d extracted",
                path.name, extracted_slides,
            )
            return

    # Fallback: extract all paragraphs if no draw:page structure found
    paragraphs = doc.getElementsByType(OdfP)
    lines: list[str] = []
    for para in paragraphs:
        text = _odf_element_text(para).strip()
        if text:
            lines.append(text)

    if not lines:
        _fail(result, "ODP presentation contains no extractable text")
        return

    _add_chunk(
        result,
        "\n".join(lines),
        source_label=f"ODP: {path.name}",
    )
    logger.debug(
        "[EXTRACTOR] ODP (fallback)  file=%s  lines=%d",
        path.name, len(lines),
    )


# ---------------------------------------------------------------------------
# Format extractors — Spreadsheets
# ---------------------------------------------------------------------------

def _extract_spreadsheet(path: Path, result: ExtractionResult) -> None:
    """Dispatch spreadsheet extraction by extension."""
    ext = path.suffix.lower()
    if ext == ".csv":
        _extract_csv(path, result)
    else:
        _extract_xlsx(path, result)


def _extract_xlsx(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .xlsx files using openpyxl.

    FR-1b Spreadsheet rule:
      - Each sheet → one ExtractedChunk labelled "Sheet: <name>".
      - Headers detected via heuristic: scan first HEADER_SCAN_ROWS
        rows and pick the row with the most non-empty cells.
      - Each data row linearized: "Key: Value | Key: Value | ..."
      - Empty rows skipped.
      - Formulas read as computed values (data_only=True).
      - CON-3: Sheets with >500 columns warned but still processed.
      - Merged cells logged at DEBUG.

    Note: openpyxl does NOT support .xls (legacy binary).
    For .xls files, _extract_xls is used instead.
    For .ods files, _extract_ods is used instead.
    """
    try:
        import openpyxl                              # type: ignore[import-untyped]
    except ImportError:
        _fail(result, "openpyxl not installed — run: pip install openpyxl")
        return

    try:
        wb = openpyxl.load_workbook(
            str(path),
            read_only  = True,
            data_only  = True,
            keep_links = False,
        )
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot open spreadsheet: {exc}")
        return

    extracted_sheets = 0

    try:
        for sheet_name in wb.sheetnames:
            try:
                ws   = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
            except Exception as exc:                 # pylint: disable=broad-except
                logger.warning(
                    "[EXTRACTOR] Sheet '%s' failed  file=%s  reason=%s",
                    sheet_name, path.name, exc,
                )
                continue

            if not rows:
                continue

            # ── Header detection heuristic ──────────────────────────────
            # Scan first HEADER_SCAN_ROWS rows; pick the one with the
            # most non-empty cells as the header row.
            best_header_row: Optional[int] = None
            best_header_count: int         = 0

            for row_idx in range(min(HEADER_SCAN_ROWS, len(rows))):
                row = rows[row_idx]
                non_empty = sum(
                    1 for c in row
                    if c is not None and str(c).strip()
                )
                if non_empty > best_header_count:
                    best_header_count = non_empty
                    best_header_row   = row_idx

            if best_header_row is None or best_header_count == 0:
                continue

            headers: list[str] = []
            for i, c in enumerate(rows[best_header_row]):
                val = str(c).strip() if c is not None else ""
                headers.append(val if val else f"Col{i + 1}")

            data_start_idx: int = best_header_row + 1

            if len(headers) > 500:
                logger.warning(
                    "[EXTRACTOR] Sheet '%s' has %d columns (CON-3)  file=%s",
                    sheet_name, len(headers), path.name,
                )

            linearized_rows: list[str] = []

            for row in rows[data_start_idx:]:
                # Log merged-cell detection
                if any(c is None for c in row):
                    non_none = sum(1 for c in row if c is not None)
                    logger.debug(
                        "[EXTRACTOR] Sheet '%s' row has %d/%d None cells "
                        "(merged?)  file=%s",
                        sheet_name, non_none, len(headers), path.name,
                    )

                pairs: list[str] = []
                for header, cell_val in zip(headers, row):
                    if cell_val is None:
                        continue
                    val = str(cell_val).strip()
                    if val:
                        pairs.append(f"{header}: {val}")
                if pairs:
                    linearized_rows.append(" | ".join(pairs))

            if not linearized_rows:
                continue

            _add_chunk(
                result,
                "\n".join(linearized_rows),
                source_label=f"Sheet: {sheet_name}",
            )
            extracted_sheets += 1

    finally:
        wb.close()

    if extracted_sheets == 0:
        _fail(result, "Spreadsheet contains no extractable data")
        return

    logger.debug(
        "[EXTRACTOR] XLSX  file=%s  sheets=%d extracted",
        path.name, extracted_sheets,
    )


def _extract_xls(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from legacy .xls files using xlrd.

    FR-1b Spreadsheet rule (XLS variant):
      - Each sheet → one ExtractedChunk labelled "Sheet: <name>".
      - First row treated as headers.
      - Each data row linearized as "Key: Value | Key: Value".
      - Empty sheets skipped.

    Note: openpyxl does NOT support .xls (binary Excel 97-2003).
    xlrd is required.
    """
    try:
        import xlrd                                  # type: ignore[import-untyped]
    except ImportError:
        _fail(result, "xlrd not installed — run: pip install xlrd")
        return

    try:
        wb = xlrd.open_workbook(str(path))
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot open XLS file: {exc}")
        return

    extracted_sheets = 0

    for sheet in wb.sheets():
        if sheet.nrows == 0:
            continue

        headers: list[str] = []
        for col in range(sheet.ncols):
            val = str(sheet.cell_value(0, col)).strip()
            headers.append(val if val else f"Col{col + 1}")

        if not any(headers):
            continue

        linearized_rows: list[str] = []

        for row_idx in range(1, sheet.nrows):
            pairs: list[str] = []
            for col_idx, header in enumerate(headers):
                val = str(sheet.cell_value(row_idx, col_idx)).strip()
                if val:
                    pairs.append(f"{header}: {val}")
            if pairs:
                linearized_rows.append(" | ".join(pairs))

        if not linearized_rows:
            continue

        _add_chunk(
            result,
            "\n".join(linearized_rows),
            source_label=f"Sheet: {sheet.name}",
        )
        extracted_sheets += 1

    if extracted_sheets == 0:
        _fail(result, "XLS file contains no extractable data")
        return

    logger.debug(
        "[EXTRACTOR] XLS  file=%s  sheets=%d extracted",
        path.name, extracted_sheets,
    )


def _extract_ods(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .ods files using pandas + odfpy engine.

    FR-1b Spreadsheet rule (ODS variant):
      - Each sheet → one ExtractedChunk labelled "Sheet: <name>".
      - Column headers used as keys.
      - Each row linearized as "Key: Value | Key: Value".
      - NaN values skipped.
      - Empty sheets skipped.

    Note: openpyxl does NOT support .ods.
    pandas with the 'odf' engine (requires odfpy) is used.
    """
    try:
        import pandas as pd                          # type: ignore[import-untyped]
    except ImportError:
        _fail(result, "pandas not installed — run: pip install pandas")
        return

    try:
        sheets = pd.read_excel(str(path), engine="odf", sheet_name=None)
    except ImportError:
        _fail(result, "odfpy not installed — run: pip install odfpy")
        return
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot open ODS file: {exc}")
        return

    extracted_sheets = 0

    for sheet_name, df in sheets.items():
        if df.empty:
            continue

        headers: list[str] = [
            str(c).strip() or f"Col{i + 1}"
            for i, c in enumerate(df.columns)
        ]

        linearized_rows: list[str] = []

        for _, row in df.iterrows():
            pairs: list[str] = []
            for header, val in zip(headers, row):
                s = str(val).strip()
                if s and s.lower() != "nan":
                    pairs.append(f"{header}: {s}")
            if pairs:
                linearized_rows.append(" | ".join(pairs))

        if not linearized_rows:
            continue

        _add_chunk(
            result,
            "\n".join(linearized_rows),
            source_label=f"Sheet: {sheet_name}",
        )
        extracted_sheets += 1

    if extracted_sheets == 0:
        _fail(result, "ODS file contains no extractable data")
        return

    logger.debug(
        "[EXTRACTOR] ODS  file=%s  sheets=%d extracted",
        path.name, extracted_sheets,
    )


def _extract_csv(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .csv files.

    FR-1b Spreadsheet rule:
      - First row treated as headers.
      - Each row linearized as "Key: Value | Key: Value".
      - Auto-detects delimiter (comma, semicolon, tab, pipe).
      - Verifies detected delimiter produces consistent column counts.
    """
    raw_text = _read_text_with_fallback(path)

    try:
        dialect = csv.Sniffer().sniff(raw_text[:4096], delimiters=",;\t|")
        # Verify: detected delimiter should produce consistent column counts
        test_reader = csv.reader(io.StringIO(raw_text[:4096]), dialect)
        test_rows   = list(test_reader)
        col_counts  = set(len(r) for r in test_rows[:5])
        if len(col_counts) > 1 or (
            len(test_rows) > 0 and max(col_counts) <= 1
        ):
            raise csv.Error("Inconsistent column count — fallback to comma")
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(raw_text), dialect)
    rows   = list(reader)

    if not rows:
        _fail(result, "CSV file is empty")
        return

    headers = [h.strip() for h in rows[0]]
    if not any(headers):
        _fail(result, "CSV has no header row")
        return

    linearized_rows: list[str] = []

    for row in rows[1:]:
        pairs: list[str] = []
        for header, value in zip(headers, row):
            val = value.strip()
            if val:
                pairs.append(f"{header}: {val}")
        if pairs:
            linearized_rows.append(" | ".join(pairs))

    if not linearized_rows:
        _fail(result, "CSV contains headers but no data rows")
        return

    _add_chunk(
        result,
        "\n".join(linearized_rows),
        source_label=f"CSV: {path.name}",
    )
    logger.debug(
        "[EXTRACTOR] CSV  file=%s  rows=%d linearized",
        path.name, len(linearized_rows),
    )


# ---------------------------------------------------------------------------
# Format extractors — Email
# ---------------------------------------------------------------------------

def _extract_eml(path: Path, result: ExtractionResult) -> None:
    """
    Extract text from .eml email files.

    FR-1b Email rule:
      - Subject, From, To, Date extracted.
      - text/plain preferred over text/html for body.
      - Attachments listed by name only (not processed — OOS-5).
      - MIME encoding (base64, quoted-printable) handled automatically.
    """
    try:
        raw_bytes = path.read_bytes()
        msg       = email_lib.message_from_bytes(
            raw_bytes,
            policy=_email_policy.default,
        )
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"Cannot parse email: {exc}")
        return

    lines: list[str] = []

    subject: str = str(msg.get("subject") or "").strip()
    sender:  str = str(msg.get("from")    or "").strip()
    to:      str = str(msg.get("to")      or "").strip()
    date:    str = str(msg.get("date")    or "").strip()

    if subject: lines.append(f"Subject: {subject}")
    if sender:  lines.append(f"From: {sender}")
    if to:      lines.append(f"To: {to}")
    if date:    lines.append(f"Date: {date}")
    lines.append("")

    body_text:   str       = ""
    attachments: list[str] = []

    if msg.is_multipart():
        for part in msg.walk():
            content_type        = part.get_content_type()
            content_disposition = str(part.get_content_disposition() or "")

            if "attachment" in content_disposition:
                fname = part.get_filename()
                if fname:
                    attachments.append(str(fname))
                continue

            if content_type == "text/plain" and not body_text:
                body_text = _safe_get_email_content(part).strip()

            elif content_type == "text/html" and not body_text:
                html_raw  = _safe_get_email_content(part)
                body_text = _strip_html(html_raw)
    else:
        body_text = _safe_get_email_content(msg).strip()

    if body_text.strip():
        lines.append(body_text.strip())

    if attachments:
        lines.append(f"[Attachments] {', '.join(attachments)}")

    if not any(line.strip() for line in lines):
        _fail(result, "Email contains no extractable content")
        return

    _add_chunk(
        result,
        "\n".join(lines),
        source_label=f"Email: {path.name}",
    )
    logger.debug(
        "[EXTRACTOR] EML  file=%s  has_attachments=%s",
        path.name, bool(attachments),
    )


# ---------------------------------------------------------------------------
# Format extractors — HTML / Web
# ---------------------------------------------------------------------------

def _extract_html(path: Path, result: ExtractionResult) -> None:
    """
    Extract semantic body text from .html / .htm files.

    FR-1b HTML/Web rule:
      - <script>, <style>, <head>, <nav>, <footer>, <aside> removed.
      - HTML entities decoded.
      - Uses built-in html.parser — no lxml dependency needed.
    """
    try:
        from bs4 import BeautifulSoup                # type: ignore[import-untyped]
    except ImportError:
        _fail(
            result,
            "beautifulsoup4 not installed — run: pip install beautifulsoup4",
        )
        return

    raw_text = _read_text_with_fallback(path)

    try:
        soup = BeautifulSoup(raw_text, "html.parser")
    except Exception as exc:                         # pylint: disable=broad-except
        _fail(result, f"HTML parse error: {exc}")
        return

    for tag in soup([
        "script", "style", "head", "nav",
        "footer", "aside", "noscript", "meta",
        "link", "button", "form", "input",
    ]):
        tag.decompose()

    raw_extracted = soup.get_text(separator="\n")
    decoded       = html_lib.unescape(raw_extracted)

    _add_chunk(result, decoded, source_label=f"HTML: {path.name}")
    logger.debug("[EXTRACTOR] HTML  file=%s  chars=%d", path.name, len(decoded))


# ---------------------------------------------------------------------------
# Shared HTML stripping utility (used by email extractor)
# ---------------------------------------------------------------------------

def _strip_html(html_text: str) -> str:
    """Strip HTML tags and return plain text. Falls back to regex."""
    try:
        from bs4 import BeautifulSoup                # type: ignore[import-untyped]
        soup = BeautifulSoup(html_text, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return html_lib.unescape(soup.get_text(separator="\n"))
    except ImportError:
        clean = re.sub(r"<[^>]+>", " ", html_text)
        return html_lib.unescape(clean)


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------
_DISPATCH: dict[str, Callable[[Path, ExtractionResult], None]] = {
    "plain_text":   _extract_plain_text,
    "source_code":  _extract_plain_text,
    "pdf":          _extract_pdf,
    "word":         _extract_docx,
    "presentation": _extract_pptx,
    "spreadsheet":  _extract_spreadsheet,
    "email":        _extract_eml,
    "web":          _extract_html,
}


# ---------------------------------------------------------------------------
# CLI entry point — standalone debug mode
# ---------------------------------------------------------------------------

def _configure_logging(verbose: bool = False) -> None:
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level   = level,
        format  = "%(asctime)s  %(levelname)-8s  %(name)s  %(message)s",
        datefmt = "%Y-%m-%d %H:%M:%S",
    )


if __name__ == "__main__":
    import argparse
    import sys

    # FIX: add project root to sys.path so 'ingestion.scanner' resolves
    # when running this file directly as a script.
    _project_root = Path(__file__).resolve().parent.parent
    sys.path.insert(0, str(_project_root))

    from ingestion.scanner import scan_directory     # noqa: E402  # type: ignore[import]

    parser = argparse.ArgumentParser(
        prog        = "extractor",
        description = "IsoCortex extractor — standalone debug mode.",
    )
    parser.add_argument("root", help="Directory to scan and extract.")
    parser.add_argument("-v", "--verbose", action="store_true")
    parser.add_argument(
        "--show-text",
        action = "store_true",
        help   = "Print 300-char preview of extracted text per chunk.",
    )
    args = parser.parse_args()

    _configure_logging(args.verbose)

    try:
        scan_result = scan_directory(args.root)
    except (FileNotFoundError, NotADirectoryError, PermissionError) as e:
        print(f"[ERROR] {e}", file=sys.stderr)
        sys.exit(1)

    if not scan_result.files:
        print("[ERROR] No supported files found.")
        sys.exit(1)

    results = extract_batch(scan_result.files)

    print()
    print("=" * 70)
    print(f"  EXTRACTION REPORT  —  {args.root}")
    print("=" * 70)

    for res in results:
        status = "OK  " if res.success else "FAIL"
        print(
            f"  [{status}]  [{res.format_category:<14}]"
            f"  chunks={res.chunk_count:<3}"
            f"  words={res.total_words:<6}"
            f"  {res.absolute_path.name}"
        )
        if not res.success:
            print(f"           error: {res.error_message}")

        if args.show_text and res.success:
            print()
            for chunk in res.chunks:
                print(f"    ── {chunk.source_label} ({chunk.word_count} words) ──")
                preview = chunk.text[:300].replace("\n", " ")
                print(f"    {preview}")
                print()

    succeeded = sum(1 for r in results if r.success)
    failed    = len(results) - succeeded

    print("=" * 70)
    print(f"  Total   : {len(results)}")
    print(f"  Success : {succeeded}")
    print(f"  Failed  : {failed}")
    print("=" * 70)
