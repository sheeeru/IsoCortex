"""
IsoCortex Desktop App — Core Engine Wrapper
=============================================
Bridges the GUI to the backend logic modules.
Handles embedding, indexing, searching, and user management
without any web server or API layer.
"""

from __future__ import annotations

import hashlib
import json
import math
import logging
import os
import re
import secrets
import shutil
import sys
import threading
import time
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Optional, TYPE_CHECKING

if TYPE_CHECKING:
    import sqlite3
    import numpy as np

try:
    import numpy as np
except ImportError:
    np = None

logger = logging.getLogger("IsoCortex.engine")

# ======================================================================
# Highlighting helper (module-level, used by engine search)
# ======================================================================

_HIGHLIGHT_MARKER = "\u0001"  # SOH char marks highlighted word boundaries

def _highlight_words(text: str, query_words: list[str]) -> str:
    """Wrap every occurrence of any query_word in *text* with marker chars."""
    if not query_words:
        return text
    # Only keep words longer than 1 char for regex matching
    words = [w for w in query_words if len(w) > 1]
    if not words:
        return text
    pattern = re.compile(
        r'\b(' + '|'.join(re.escape(w) for w in words) + r')\b',
        re.IGNORECASE,
    )
    # Use a lambda to avoid regex replacement string escape issues
    return pattern.sub(lambda m: _HIGHLIGHT_MARKER + m.group(1) + _HIGHLIGHT_MARKER, text)

# Public alias for use by other modules (e.g. search screen highlighting)
HIGHLIGHT_MARKER = _HIGHLIGHT_MARKER

# ======================================================================
# Configuration
# ======================================================================

DEFAULT_DATA_DIR = Path.home() / ".isocortex"
DEFAULT_MODEL_NAME = "all-MiniLM-L6-v2 (ONNX)"
DEFAULT_VECTOR_DIM = 384
EMBEDDING_MODEL_DIR = DEFAULT_DATA_DIR / "embedding_model"


def embedding_model_exists() -> bool:
    """Check if the ONNX embedding model is already downloaded."""
    onnx_file = EMBEDDING_MODEL_DIR / "model.onnx"
    tok_file = EMBEDDING_MODEL_DIR / "tokenizer.json"
    return onnx_file.exists() and tok_file.exists()


def download_embedding_model(
    on_progress: Optional[Callable[[str], None]] = None,
) -> Path:
    """
    Download the all-MiniLM-L6-v2 ONNX model from HuggingFace.

    Downloads:
        - model.onnx  (the ONNX embedding model, ~90 MB)
        - tokenizer.json  (the HuggingFace tokenizer)

    Args:
        on_progress: Callback(status_text) for UI feedback.

    Returns:
        Path to the embedding model directory.
    """
    EMBEDDING_MODEL_DIR.mkdir(parents=True, exist_ok=True)

    if on_progress:
        on_progress("Checking for embedding model...")

    # If already present, skip
    if embedding_model_exists():
        if on_progress:
            on_progress("Embedding model already present.")
        return EMBEDDING_MODEL_DIR

    try:
        from huggingface_hub import hf_hub_download
    except ImportError:
        raise RuntimeError(
            "huggingface_hub is required for model download. "
            "Install with: pip install huggingface_hub"
        )

    repo_id = "sentence-transformers/all-MiniLM-L6-v2"

    # Download ONNX model
    if on_progress:
        on_progress("Downloading embedding model (model.onnx)...")
    onnx_path = hf_hub_download(
        repo_id=repo_id,
        filename="onnx/model.onnx",
        local_dir=str(EMBEDDING_MODEL_DIR),
        local_dir_use_symlinks=False,  # type: ignore[arg-type]
    )
    # Move from subfolder if needed
    downloaded_onnx = Path(onnx_path)
    target_onnx = EMBEDDING_MODEL_DIR / "model.onnx"
    if downloaded_onnx != target_onnx and not target_onnx.exists():
        downloaded_onnx.rename(target_onnx)

    # Download tokenizer
    if on_progress:
        on_progress("Downloading tokenizer...")
    tok_path = hf_hub_download(
        repo_id=repo_id,
        filename="tokenizer.json",
        local_dir=str(EMBEDDING_MODEL_DIR),
        local_dir_use_symlinks=False,  # type: ignore[arg-type]
    )
    downloaded_tok = Path(tok_path)
    target_tok = EMBEDDING_MODEL_DIR / "tokenizer.json"
    if downloaded_tok != target_tok and not target_tok.exists():
        downloaded_tok.rename(target_tok)

    # Cleanup empty subdirectories created by hf_hub_download
    for sub in (EMBEDDING_MODEL_DIR / "onnx",):
        if sub.is_dir():
            try:
                sub.rmdir()
            except OSError:
                pass

    logger.info("Embedding model downloaded to: %s", EMBEDDING_MODEL_DIR)
    if on_progress:
        on_progress("Embedding model ready.")
    return EMBEDDING_MODEL_DIR


DEFAULT_HNSW_M = 16
DEFAULT_HNSW_EF_CONSTRUCTION = 200
DEFAULT_HNSW_EF_SEARCH = 50
DEFAULT_CHUNK_SIZE = 350
DEFAULT_OVERLAP = 38
DEFAULT_BATCH_SIZE = 64
DEFAULT_INDEX_NAME = "default"

MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB max per file

# Supported file extensions for scanning
SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".log", ".rst", ".cfg", ".ini", ".py", ".cpp", ".c",
    ".h", ".js", ".ts", ".java", ".go", ".rs", ".rb", ".pdf", ".docx",
    ".odt", ".pptx", ".odp", ".xlsx", ".xls", ".ods", ".csv", ".json",
    ".tsv", ".eml", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp",
}


# ======================================================================
# Data Classes
# ======================================================================

@dataclass
class IndexInfo:
    """Lightweight index info for display."""
    name: str
    description: str
    vector_count: int
    deleted_count: int
    created_at: str
    updated_at: str
    index_size_mb: float = 0.0
    healthy: bool = True

    def to_dict(self) -> dict:
        return {
            "name": self.name,
            "description": self.description,
            "vector_count": self.vector_count,
            "deleted_count": self.deleted_count,
            "created_at": self.created_at,
            "updated_at": self.updated_at,
            "index_size_mb": round(self.index_size_mb, 2),
            "healthy": self.healthy,
        }


@dataclass
class SearchResult:
    """A single semantic search result."""
    rank: int
    text: str
    source_file: str
    source_label: str
    score: float
    chunk_index: int
    format_category: str = ""
    page_number: int = 0
    full_text: str = ""
    matched_segment: str = ""


@dataclass
class IngestionStats:
    """Statistics from a file ingestion operation."""
    files_scanned: int = 0
    files_accepted: int = 0
    files_processed: int = 0
    files_failed: int = 0
    total_chunks: int = 0
    total_vectors: int = 0
    total_tokens: int = 0
    elapsed_seconds: float = 0.0


# ======================================================================
# Core Engine
# ======================================================================

class IsoCortexEngine:
    """
    Central engine that wraps all backend logic for the desktop GUI.

    Handles:
      - User authentication (bcrypt + local JWT)
      - Embedding model management
      - Index lifecycle (create, list, load, delete, search)
      - File ingestion pipeline (scan -> extract -> chunk -> embed -> index)
      - Settings persistence
    """

    def __init__(self, data_dir: Optional[Path] = None):
        self._data_dir = data_dir or DEFAULT_DATA_DIR
        self._data_dir.mkdir(parents=True, exist_ok=True)
        self._indices_dir = self._data_dir / "indices"
        self._indices_dir.mkdir(parents=True, exist_ok=True)
        self._db_path = self._data_dir / "isocortex.db"

        # In-memory state
        self._loaded_indices: dict[str, dict] = {}  # name -> {vectors, metadata}
        self._embed_model = None
        self._tokenizer = None
        self._ort_session = None
        self._current_user: Optional[dict] = None
        self._jwt_secret = self._get_or_create_jwt_secret()
        self._lock = threading.Lock()
        self._db_conn: sqlite3.Connection | None = None  # type: ignore[name-defined]  # Reusable DB connection (main thread)
        self._thread_db_connections: dict = {}  # Thread-local DB connections

        # Plugin system
        self._plugin_manager = None

        # Initialize database
        self._init_database()

        # Load plugins
        try:
            from desktop_app.plugins import PluginManager
            self._plugin_manager = PluginManager()
            count = self._plugin_manager.load_all()
            if count:
                logger.info("Loaded %d plugins", count)
        except Exception as exc:
            logger.debug("Plugin system init failed: %s", exc)

        logger.info(
            "Engine initialized  data_dir=%s",
            self._data_dir,
        )

    # ------------------------------------------------------------------
    # Properties
    # ------------------------------------------------------------------

    @property
    def data_dir(self) -> Path:
        return self._data_dir

    @property
    def indices_dir(self) -> Path:
        return self._indices_dir

    @property
    def is_authenticated(self) -> bool:
        with self._lock:
            return self._current_user is not None

    @property
    def current_user(self) -> Optional[dict]:
        with self._lock:
            return self._current_user

    @property
    def plugin_manager(self):
        """Return the plugin manager instance."""
        return self._plugin_manager

    # ------------------------------------------------------------------
    # JWT Secret Management
    # ------------------------------------------------------------------

    def _get_or_create_jwt_secret(self) -> str:
        secret_path = self._data_dir / ".jwt_secret"
        if secret_path.exists():
            try:
                return secret_path.read_text(encoding="utf-8").strip()
            except Exception:
                pass
        secret = secrets.token_urlsafe(64)
        try:
            secret_path.write_text(secret, encoding="utf-8")
            os.chmod(secret_path, 0o600)
        except Exception as exc:
            logger.warning("Failed to persist JWT secret: %s", exc)
        return secret

    # ------------------------------------------------------------------
    # Database
    # ------------------------------------------------------------------

    def _get_db(self) -> sqlite3.Connection:  # type: ignore[name-defined]
        """Get a SQLite database connection.

        For the main thread, reuses a single cached connection.
        For worker threads, creates (or reuses from thread-local) a separate connection
        to avoid SQLite threading issues.
        """
        import sqlite3
        thread_id = threading.get_ident()

        # Main thread uses the cached connection (protected by _lock)
        if thread_id == threading.main_thread().ident:
            with self._lock:
                if self._db_conn is not None:
                    try:
                        self._db_conn.execute("SELECT 1")
                        return self._db_conn
                    except Exception:
                        self._db_conn = None

                self._db_conn = sqlite3.connect(
                    str(self._db_path),
                    check_same_thread=False,
                    timeout=5.0,
                )
                self._db_conn.row_factory = sqlite3.Row
                self._db_conn.execute("PRAGMA journal_mode=WAL")
                self._db_conn.execute("PRAGMA foreign_keys=ON")
                return self._db_conn

        # Worker threads: use a thread-local connection
        if thread_id not in self._thread_db_connections:
            conn = sqlite3.connect(
                str(self._db_path),
                check_same_thread=True,
                timeout=10.0,
            )
            conn.row_factory = sqlite3.Row
            conn.execute("PRAGMA journal_mode=WAL")
            conn.execute("PRAGMA foreign_keys=ON")
            self._thread_db_connections[thread_id] = conn
        return self._thread_db_connections[thread_id]

    def get_db_connection(self) -> sqlite3.Connection:  # type: ignore[name-defined]
        """Public wrapper around _get_db().

        Returns a thread-safe SQLite connection with ``row_factory=sqlite3.Row``.
        Callers should NOT cache the returned connection across threads.
        """
        return self._get_db()

    def _init_database(self) -> None:
        """Initialize all required database tables."""
        conn = self._get_db()
        conn.executescript("""
            CREATE TABLE IF NOT EXISTS users (
                user_id     TEXT PRIMARY KEY,
                username    TEXT NOT NULL UNIQUE,
                email       TEXT NOT NULL UNIQUE,
                password_hash TEXT NOT NULL,
                role        TEXT NOT NULL DEFAULT 'user',
                is_active   INTEGER NOT NULL DEFAULT 1,
                created_at  TEXT NOT NULL DEFAULT (datetime('now')),
                updated_at  TEXT NOT NULL DEFAULT (datetime('now')),
                last_login  TEXT,
                locked_at   TEXT,
                failed_login_attempts INTEGER NOT NULL DEFAULT 0
            );

            CREATE TABLE IF NOT EXISTS analytics (
                event_id   TEXT PRIMARY KEY,
                event_type TEXT NOT NULL,
                user_id    TEXT,
                metadata   TEXT NOT NULL DEFAULT '{}',
                created_at TEXT NOT NULL DEFAULT (datetime('now'))
            );

            CREATE TABLE IF NOT EXISTS documents (
                doc_id      TEXT PRIMARY KEY,
                index_name  TEXT NOT NULL,
                file_path   TEXT NOT NULL,
                file_hash   TEXT NOT NULL DEFAULT '',
                format_category TEXT NOT NULL DEFAULT '',
                chunk_count INTEGER NOT NULL DEFAULT 0,
                word_count  INTEGER NOT NULL DEFAULT 0,
                status      TEXT NOT NULL DEFAULT 'indexed',
                created_at  TEXT NOT NULL DEFAULT (datetime('now'))
            );

            CREATE TABLE IF NOT EXISTS chat_conversations (
                conversation_id TEXT PRIMARY KEY,
                user_id        TEXT,
                title           TEXT NOT NULL DEFAULT 'New Chat',
                created_at      TEXT NOT NULL DEFAULT (datetime('now')),
                updated_at      TEXT NOT NULL DEFAULT (datetime('now'))
            );

            CREATE TABLE IF NOT EXISTS chat_messages (
                message_id     TEXT PRIMARY KEY,
                conversation_id TEXT NOT NULL,
                role           TEXT NOT NULL DEFAULT 'user',
                content        TEXT NOT NULL DEFAULT '',
                sources_json   TEXT NOT NULL DEFAULT '[]',
                token_count    INTEGER NOT NULL DEFAULT 0,
                created_at     TEXT NOT NULL DEFAULT (datetime('now')),
                FOREIGN KEY (conversation_id) REFERENCES chat_conversations(conversation_id)
                    ON DELETE CASCADE
            );

            CREATE TABLE IF NOT EXISTS watched_folders (
                folder_id   TEXT PRIMARY KEY,
                folder_path TEXT NOT NULL UNIQUE,
                index_name  TEXT NOT NULL DEFAULT 'default',
                is_active   INTEGER NOT NULL DEFAULT 1,
                created_at  TEXT NOT NULL DEFAULT (datetime('now'))
            );

            CREATE TABLE IF NOT EXISTS search_history (
                history_id  TEXT PRIMARY KEY,
                query       TEXT NOT NULL,
                results_count INTEGER NOT NULL DEFAULT 0,
                created_at  TEXT NOT NULL DEFAULT (datetime('now'))
            );

            CREATE TABLE IF NOT EXISTS bookmarks (
                bookmark_id TEXT PRIMARY KEY,
                history_id  TEXT,
                query       TEXT NOT NULL,
                file_path   TEXT NOT NULL DEFAULT '',
                chunk_text  TEXT NOT NULL DEFAULT '',
                score       REAL NOT NULL DEFAULT 0,
                note        TEXT NOT NULL DEFAULT '',
                created_at  TEXT NOT NULL DEFAULT (datetime('now'))
            );
        """)
        conn.commit()
        logger.info("Database initialized")

        # Performance indexes for common queries
        conn.executescript("""
            CREATE INDEX IF NOT EXISTS idx_documents_idx_hash ON documents(index_name, file_hash);
            CREATE INDEX IF NOT EXISTS idx_chat_messages_conv ON chat_messages(conversation_id);
            CREATE INDEX IF NOT EXISTS idx_analytics_type ON analytics(event_type);
            CREATE INDEX IF NOT EXISTS idx_search_history_time ON search_history(created_at DESC);
        """)

        # Secure database file permissions
        try:
            os.chmod(str(self._db_path), 0o600)
        except OSError:
            pass

    # ------------------------------------------------------------------
    # User Management
    # ------------------------------------------------------------------

    def _hash_password(self, password: str) -> str:
        """Hash a password using bcrypt (preferred) or PBKDF2-SHA256 (fallback)."""
        try:
            import bcrypt
            return bcrypt.hashpw(
                password.encode("utf-8"),
                bcrypt.gensalt(rounds=12),
            ).decode("utf-8")
        except ImportError:
            logger.warning(
                "bcrypt not installed — falling back to PBKDF2-SHA256. "
                "Install bcrypt for stronger security: pip install bcrypt"
            )
            salt = os.urandom(32)
            derived = hashlib.pbkdf2_hmac(
                "sha256", password.encode("utf-8"), salt, iterations=600_000,
            )
            return f"pbkdf2${salt.hex()}${derived.hex()}"

    def _verify_password(self, password: str, password_hash: str) -> bool:
        """Verify a password against its hash."""
        try:
            import bcrypt
            return bcrypt.checkpw(
                password.encode("utf-8"),
                password_hash.encode("utf-8"),
            )
        except ImportError:
            pass
        # PBKDF2 fallback
        if password_hash.startswith("pbkdf2$"):
            parts = password_hash.split("$")
            if len(parts) != 3:
                return False
            salt = bytes.fromhex(parts[1])
            stored = parts[2]
            derived = hashlib.pbkdf2_hmac(
                "sha256", password.encode("utf-8"), salt, iterations=600_000,
            )
            return secrets.compare_digest(derived.hex(), stored)
        # Legacy sha256$ format (v1.0 compat — still verify but log warning)
        if password_hash.startswith("sha256$"):
            logger.warning("Legacy sha256 password hash detected — please re-hash with bcrypt")
            parts = password_hash.split("$")
            if len(parts) != 3:
                return False
            salt, stored_hash = parts[1], parts[2]
            computed = hashlib.sha256(
                (salt + password).encode("utf-8")
            ).hexdigest()
            return secrets.compare_digest(computed, stored_hash)
        return False

    def create_user(self, username: str, email: str, password: str) -> dict:
        """Create a new user account."""
        if not username or not email or not password:
            raise ValueError("Username, email, and password are required")
        if len(password) < 8:
            raise ValueError("Password must be at least 8 characters")

        # Password complexity requirements
        if not any(c.isupper() for c in password):
            raise ValueError("Password must contain at least one uppercase letter")
        if not any(c.islower() for c in password):
            raise ValueError("Password must contain at least one lowercase letter")
        if not any(c.isdigit() for c in password):
            raise ValueError("Password must contain at least one digit")
        if not any(c in "!@#$%^&*()_+-=[]{}|;:',.<>?/`~" for c in password):
            raise ValueError("Password must contain at least one special character")
        if len(password) > 128:
            raise ValueError("Password must be at most 128 characters")

        # Input validation
        if len(username) < 3 or len(username) > 32:
            raise ValueError("Username must be between 3 and 32 characters")
        if not all(c.isalnum() or c in "-_." for c in username):
            raise ValueError("Username can only contain letters, numbers, hyphens, underscores, and dots")
        # Email validation (basic)
        email_parts = email.split("@")
        if len(email_parts) != 2 or not email_parts[0] or not email_parts[1] or "." not in email_parts[1]:
            raise ValueError("Invalid email address format")

        user_id = str(uuid.uuid4())
        pwd_hash = self._hash_password(password)

        conn = self._get_db()
        try:
            conn.execute(
                """INSERT INTO users (user_id, username, email, password_hash)
                   VALUES (?, ?, ?, ?)""",
                (user_id, username, email, pwd_hash),
            )
            conn.commit()
        except Exception as exc:
            if "UNIQUE constraint" in str(exc):
                raise ValueError(
                    f"Username '{username}' or email '{email}' already exists"
                ) from exc
            raise

        logger.info("User created  username=%s  user_id=%s", username, user_id)
        row = conn.execute(
            "SELECT user_id, username, email, role FROM users WHERE username = ?",
            (username,),
        ).fetchone()
        return {
            "user_id": row["user_id"],
            "username": row["username"],
            "email": row["email"],
            "role": row["role"],
        }

    def authenticate(self, username_or_email: str, password: str) -> dict:
        """Authenticate a user and return session data."""
        if not username_or_email or not password:
            raise ValueError("Username and password are required")

        conn = self._get_db()
        row = conn.execute(
            """SELECT user_id, username, email, password_hash, role,
                      is_active, locked_at, failed_login_attempts
               FROM users
               WHERE (username = ? OR email = ?) AND is_active = 1""",
            (username_or_email, username_or_email),
        ).fetchone()

        if row is None:
            raise ValueError("Invalid username or password")

        # Check account lockout (5 failed attempts -> 30 min lock)
        if row["failed_login_attempts"] >= 5 and row["locked_at"]:
            try:
                locked_until = datetime.fromisoformat(row["locked_at"])
                lock_duration = locked_until.timestamp() + 1800
                if datetime.now(timezone.utc).timestamp() < lock_duration:
                    remaining = int(lock_duration - datetime.now(timezone.utc).timestamp())
                    raise ValueError(
                        f"Account locked. Try again in {remaining // 60} minutes"
                    )
                else:
                    # Lock expired, reset
                    conn.execute(
                        "UPDATE users SET failed_login_attempts = 0, locked_at = NULL WHERE user_id = ?",
                        (row["user_id"],),
                    )
                    conn.commit()
            except ValueError:
                raise  # Re-raise the lockout ValueError
            except (TypeError, AttributeError, OverflowError):
                pass

        if not self._verify_password(password, row["password_hash"]):
            # Increment failed attempts
            new_attempts = row["failed_login_attempts"] + 1
            lock_at = datetime.now(timezone.utc).isoformat() if new_attempts >= 5 else None
            conn.execute(
                "UPDATE users SET failed_login_attempts = ?, locked_at = ? WHERE user_id = ?",
                (new_attempts, lock_at, row["user_id"]),
            )
            conn.commit()
            raise ValueError("Invalid username or password")

        # Reset failed attempts on successful login
        now = datetime.now(timezone.utc).isoformat()
        conn.execute(
            "UPDATE users SET failed_login_attempts = 0, locked_at = NULL, last_login = ? WHERE user_id = ?",
            (now, row["user_id"]),
        )
        conn.commit()

        user = {
            "user_id": row["user_id"],
            "username": row["username"],
            "email": row["email"],
            "role": row["role"],
        }
        with self._lock:
            self._current_user = user

        # Record analytics
        self._record_event("user_login", metadata={"username": row["username"]})

        logger.info("User authenticated  username=%s", row["username"])
        return user

    def logout(self) -> None:
        """Log out the current user."""
        with self._lock:
            self._current_user = None
        logger.info("User logged out")

    def change_password(self, current_password: str, new_password: str) -> None:
        """Change the current user's password."""
        with self._lock:
            user = self._current_user
        if not user:
            raise ValueError("No user is currently authenticated")

        conn = self._get_db()
        row = conn.execute(
            "SELECT password_hash FROM users WHERE user_id = ?",
            (user["user_id"],),
        ).fetchone()

        if row is None:
            raise ValueError("User not found")

        if not self._verify_password(current_password, row["password_hash"]):
            raise ValueError("Current password is incorrect")

        # Validate new password complexity
        if len(new_password) < 8:
            raise ValueError("Password must be at least 8 characters")
        if len(new_password) > 128:
            raise ValueError("Password must be at most 128 characters")
        if not any(c.isupper() for c in new_password):
            raise ValueError("Password must contain at least one uppercase letter")
        if not any(c.islower() for c in new_password):
            raise ValueError("Password must contain at least one lowercase letter")
        if not any(c.isdigit() for c in new_password):
            raise ValueError("Password must contain at least one digit")
        if not any(c in "!@#$%^&*()_+-=[]{}|;:',.<>?/`~" for c in new_password):
            raise ValueError("Password must contain at least one special character")
        if new_password == current_password:
            raise ValueError("New password must be different from current password")

        new_hash = self._hash_password(new_password)
        conn.execute(
            "UPDATE users SET password_hash = ?, updated_at = datetime('now') WHERE user_id = ?",
            (new_hash, user["user_id"]),
        )
        conn.commit()

        self._record_event("password_changed", metadata={"user_id": user["user_id"]})
        logger.info("Password changed for user %s", user["username"])

    def is_first_run(self) -> bool:
        """Check if this is the first run (no users exist)."""
        conn = self._get_db()
        count = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
        return count == 0

    # ------------------------------------------------------------------
    # Analytics
    # ------------------------------------------------------------------

    def _record_event(self, event_type: str, metadata: dict | None = None) -> None:
        """Record an analytics event."""
        event_id = str(uuid.uuid4())
        meta = json.dumps(metadata or {}, default=str)
        user_id = self._current_user["user_id"] if self._current_user else None

        conn = self._get_db()
        try:
            conn.execute(
                "INSERT INTO analytics (event_id, event_type, user_id, metadata) VALUES (?, ?, ?, ?)",
                (event_id, event_type, user_id, meta),
            )
            conn.commit()
        except Exception as exc:
            logger.warning("Failed to record event: %s", exc)

    def get_system_stats(self) -> dict:
        """Get system statistics for the dashboard."""
        conn = self._get_db()
        index_count = len(self._list_index_dirs())

        doc_row = conn.execute("SELECT COUNT(*) FROM documents").fetchone()
        total_documents = doc_row[0] if doc_row else 0

        total_vectors = 0
        # Count vectors from all loaded indexes in memory
        with self._lock:
            for name, idx in self._loaded_indices.items():
                if idx.get("vectors") is not None:
                    total_vectors += idx["vectors"].shape[0]
            loaded_names = set(self._loaded_indices.keys())
        for idx_name in self._list_index_dirs():
            if idx_name not in loaded_names:
                info = self._read_index_info(idx_name)
                if info:
                    total_vectors += info.get("vector_count", 0)

        search_row = conn.execute(
            "SELECT COUNT(*) FROM analytics WHERE event_type = 'search'"
        ).fetchone()
        total_searches = search_row[0] if search_row else 0

        # Data directory size (cached for 30 seconds to avoid blocking)
        total_size_mb = 0.0
        if self._data_dir.exists():
            now = time.time()
            if (not hasattr(self, '_data_size_cache')
                    or not isinstance(self._data_size_cache, (list, tuple))
                    or len(self._data_size_cache) < 2
                    or now - self._data_size_cache[0] > 30):
                size_bytes = 0
                for f in self._data_dir.rglob("*"):
                    if f.is_file():
                        try:
                            size_bytes += f.stat().st_size
                        except OSError:
                            pass
                self._data_size_cache = (now, size_bytes)
            total_size_mb = self._data_size_cache[1] / (1024 * 1024)

        return {
            "index_count": index_count,
            "total_documents": total_documents,
            "total_vectors": total_vectors,
            "total_searches": total_searches,
            "data_size_mb": round(total_size_mb, 1),
        }

    # ------------------------------------------------------------------
    # Embedding Model
    # ------------------------------------------------------------------

    def _get_model_dir(self) -> Path:
        """Return the directory containing the ONNX model and tokenizer files."""
        # 1. Auto-downloaded model in ~/.isocortex/embedding_model/
        if (EMBEDDING_MODEL_DIR / "model.onnx").exists() and \
           (EMBEDDING_MODEL_DIR / "tokenizer.json").exists():
            return EMBEDDING_MODEL_DIR

        # 2. Bundled model next to this file
        bundled = Path(__file__).parent / "assets" / "model"
        if bundled.exists() and (bundled / "model.onnx").exists():
            return bundled

        # 3. PyInstaller _MEIPASS
        if hasattr(sys, "_MEIPASS"):
            bundled = Path(sys._MEIPASS) / "desktop_app" / "assets" / "model"  # type: ignore[attr-defined]
            if bundled.exists() and (bundled / "model.onnx").exists():
                return bundled

        raise FileNotFoundError(
            "ONNX model files not found. Expected model.onnx and tokenizer.json "
            "in ~/.isocortex/embedding_model/ or desktop_app/assets/model/"
        )

    def ensure_model(self) -> bool:
        """Ensure the embedding model (ONNX) is loaded. Returns True on success."""
        if self._embed_model is not None:
            return True

        if np is None:
            logger.error(
                "numpy is not installed but is required for embeddings. "
                "Install with: pip install numpy"
            )
            return False

        try:
            import onnxruntime as ort
            from tokenizers import Tokenizer

            model_dir = self._get_model_dir()
            onnx_path = model_dir / "model.onnx"
            tok_path = model_dir / "tokenizer.json"

            logger.info("Loading ONNX embedding model from: %s", onnx_path)

            # Load tokenizer
            self._tokenizer = Tokenizer.from_file(str(tok_path))
            self._tokenizer.enable_truncation(max_length=128)
            self._tokenizer.enable_padding(length=128)

            # Load ONNX model
            sess_options = ort.SessionOptions()
            sess_options.graph_optimization_level = ort.GraphOptimizationLevel.ORT_ENABLE_ALL
            self._ort_session = ort.InferenceSession(
                str(onnx_path),
                sess_options=sess_options,
            )

            # Verify output dimensions
            test_encoding = self._tokenizer.encode("test")
            ort_inputs = {
                "input_ids": np.array([test_encoding.ids], dtype=np.int64),
                "attention_mask": np.array([test_encoding.attention_mask], dtype=np.int64),
                "token_type_ids": np.array([test_encoding.type_ids], dtype=np.int64),
            }
            test_output = self._ort_session.run(None, ort_inputs)
            # Mean pooling over token embeddings (CLS-style pooling)
            last_hidden = test_output[0]  # (1, seq_len, hidden_dim)
            mask_expanded = ort_inputs["attention_mask"][:, :, np.newaxis]  # (1, seq_len, 1)
            sum_emb = (last_hidden * mask_expanded).sum(axis=1)
            count = mask_expanded.sum(axis=1).clip(min=1e-9)
            mean_pooled = (sum_emb / count)  # (1, hidden_dim)
            # L2 normalize
            norms = np.linalg.norm(mean_pooled, axis=1, keepdims=True)
            normalized = mean_pooled / norms.clip(min=1e-9)

            if normalized.shape[1] != DEFAULT_VECTOR_DIM:
                logger.error(
                    "Model output dimension %d != expected %d",
                    normalized.shape[1], DEFAULT_VECTOR_DIM,
                )
                self._embed_model = False
                return False

            self._embed_model = True  # Model loaded successfully
            logger.info("ONNX embedding model loaded successfully (dim=%d)", normalized.shape[1])
            return True

        except ImportError as exc:
            logger.error(
                "Missing dependency for ONNX model: %s. "
                "Install with: pip install onnxruntime tokenizers",
                exc,
            )
            return False
        except FileNotFoundError as exc:
            logger.error("Model files not found: %s", exc)
            return False
        except Exception as exc:
            logger.error("Failed to load ONNX embedding model: %s", exc, exc_info=True)
            return False

    def _encode_texts(self, texts: list[str]) -> Optional[np.ndarray]:  # type: ignore[name-defined]
        """Run texts through ONNX model and return L2-normalized embeddings."""
        if not texts or self._embed_model is None:
            return None

        if np is None:
            logger.error("numpy is required for text encoding")
            return None

        assert np is not None  # for Pylance

        # Tokenize all texts
        encodings = self._tokenizer.encode_batch(texts)

        input_ids = np.array([e.ids for e in encodings], dtype=np.int64)
        attention_mask = np.array([e.attention_mask for e in encodings], dtype=np.int64)
        token_type_ids = np.array([e.type_ids for e in encodings], dtype=np.int64)

        # Check if model needs token_type_ids
        input_names = {inp.name for inp in self._ort_session.get_inputs()}
        ort_inputs = {
            "input_ids": input_ids,
            "attention_mask": attention_mask,
        }
        if "token_type_ids" in input_names:
            ort_inputs["token_type_ids"] = token_type_ids

        # Run inference (thread-safe: ONNX InferenceSession is not safe
        # for concurrent use on a single session, so we serialize access)
        with self._lock:
            outputs = self._ort_session.run(None, ort_inputs)
        last_hidden = outputs[0]  # (batch, seq_len, hidden_dim)

        # Mean pooling (CLS-style pooling)
        mask_expanded = attention_mask[:, :, np.newaxis].astype(np.float32)
        sum_emb = (last_hidden * mask_expanded).sum(axis=1)
        count = mask_expanded.sum(axis=1).clip(min=1e-9)
        mean_pooled = sum_emb / count

        # L2 normalize
        norms = np.linalg.norm(mean_pooled, axis=1, keepdims=True)
        normalized = mean_pooled / norms.clip(min=1e-9)

        return normalized.astype(np.float32)

    def embed_text(self, text: str) -> Optional[np.ndarray]:  # type: ignore[name-defined]
        """Embed a single text string into a 384-dim vector."""
        if not self.ensure_model():
            return None
        try:
            vectors = self._encode_texts([text])
            if vectors is not None:
                return vectors[0]
            return None
        except Exception as exc:
            logger.error("Embedding failed: %s", exc)
            return None

    def embed_batch(self, texts: list[str]) -> Optional[np.ndarray]:  # type: ignore[name-defined]
        """Embed multiple texts into vectors."""
        if not texts or not self.ensure_model():
            return None
        try:
            return self._encode_texts(texts)
        except Exception as exc:
            logger.error("Batch embedding failed: %s", exc)
            return None

    def get_model_status(self) -> dict:
        """Return status of the embedding model."""
        if self._embed_model is not None:
            return {
                "loaded": True,
                "model_name": DEFAULT_MODEL_NAME,
                "device": "ONNX Runtime",
                "dimension": DEFAULT_VECTOR_DIM,
            }
        return {
            "loaded": False,
            "model_name": DEFAULT_MODEL_NAME,
            "device": "N/A",
            "dimension": DEFAULT_VECTOR_DIM,
        }

    # ------------------------------------------------------------------
    # Index Management
    # ------------------------------------------------------------------

    def ensure_default_index(self) -> str:
        """Ensure the default index exists. Create it if missing. Returns the index name."""
        if self._read_index_info(DEFAULT_INDEX_NAME) is not None:
            return DEFAULT_INDEX_NAME
        self.create_index(DEFAULT_INDEX_NAME, "Default index — all uploaded files go here")
        logger.info("Auto-created default index")
        return DEFAULT_INDEX_NAME

    def _list_index_dirs(self) -> list[str]:
        """List all index directory names."""
        if not self._indices_dir.exists():
            return []
        return [
            d.name for d in self._indices_dir.iterdir()
            if d.is_dir() and (d / "index_info.json").exists()
        ]

    def _read_index_info(self, name: str) -> Optional[dict]:
        """Read index_info.json from disk."""
        path = self._indices_dir / name / "index_info.json"
        if not path.exists():
            return None
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return None

    def _write_index_info(self, name: str, info: dict) -> None:
        """Write index_info.json atomically."""
        idx_dir = self._indices_dir / name
        idx_dir.mkdir(parents=True, exist_ok=True)
        target = idx_dir / "index_info.json"
        tmp = target.with_suffix(".tmp")
        tmp.write_text(json.dumps(info, indent=2, default=str), encoding="utf-8")
        tmp.replace(target)

    def _safe_index_name(self, name: str) -> str:
        """Sanitize index name to prevent path traversal."""
        return "".join(c for c in name if c.isalnum() or c in "-_").strip()

    def list_indexes(self) -> list[IndexInfo]:
        """List all indexes with their metadata."""
        result = []
        for name in self._list_index_dirs():
            info = self._read_index_info(name)
            if info is None:
                continue

            # Calculate size
            idx_dir = self._indices_dir / name
            total_size = 0
            try:
                total_size = sum(
                    f.stat().st_size for f in idx_dir.iterdir() if f.is_file()
                )
            except OSError:
                pass

            result.append(IndexInfo(
                name=name,
                description=info.get("description", ""),
                vector_count=info.get("vector_count", 0),
                deleted_count=info.get("deleted_count", 0),
                created_at=info.get("created_at", ""),
                updated_at=info.get("updated_at", ""),
                index_size_mb=total_size / (1024 * 1024),
            ))
        return result

    def create_index(self, name: str, description: str = "") -> str:
        """Create a new empty index."""
        safe_name = self._safe_index_name(name)
        if not safe_name:
            raise ValueError("Invalid index name")

        idx_dir = self._indices_dir / safe_name
        if idx_dir.exists() and (idx_dir / "index_info.json").exists():
            raise FileExistsError(f"Index '{name}' already exists")

        now = datetime.now(timezone.utc).isoformat()
        info = {
            "name": safe_name,
            "description": description,
            "created_at": now,
            "updated_at": now,
            "embedding_model": DEFAULT_MODEL_NAME,
            "vector_count": 0,
            "deleted_count": 0,
            "hnsw_params": {
                "M": DEFAULT_HNSW_M,
                "ef_construction": DEFAULT_HNSW_EF_CONSTRUCTION,
                "ef_search": DEFAULT_HNSW_EF_SEARCH,
                "metric": "cosine",
            },
            "chunk_config": {
                "chunk_size": DEFAULT_CHUNK_SIZE,
                "overlap": DEFAULT_OVERLAP,
            },
        }

        idx_dir.mkdir(parents=True, exist_ok=True)
        self._write_index_info(safe_name, info)

        self._record_event("index_created", metadata={"index_name": safe_name})
        logger.info("Index created  name=%s", safe_name)
        return safe_name

    def delete_index(self, name: str) -> None:
        """Delete an index from disk and memory."""
        safe_name = self._safe_index_name(name)
        idx_dir = self._indices_dir / safe_name
        if not idx_dir.exists():
            raise FileNotFoundError(f"Index '{name}' not found")

        # Unload from memory
        with self._lock:
            self._loaded_indices.pop(safe_name, None)

        # Remove from disk
        shutil.rmtree(idx_dir, ignore_errors=True)

        # Remove document records from DB
        conn = self._get_db()
        try:
            conn.execute("DELETE FROM documents WHERE index_name = ?", (safe_name,))
            conn.commit()
        except Exception:
            pass

        self._record_event("index_deleted", metadata={"index_name": safe_name})
        logger.info("Index deleted  name=%s", safe_name)

    def load_index(self, name: str) -> bool:
        """Load an index from disk into memory."""
        safe_name = self._safe_index_name(name)
        with self._lock:
            if safe_name in self._loaded_indices:
                return True

        idx_dir = self._indices_dir / safe_name
        vectors_path = idx_dir / "vectors.bin"
        metadata_path = idx_dir / "metadata.json"

        vectors = None
        metadata = []

        if np is None:
            logger.error("numpy is required for loading index vectors")
            with self._lock:
                self._loaded_indices[safe_name] = {
                    "vectors": None,
                    "metadata": [],
                }
            return False

        if vectors_path.exists():
            try:
                import struct
                file_size = vectors_path.stat().st_size
                if file_size >= 20:
                    with open(str(vectors_path), "rb") as f:
                        # Read 5 x uint32 header (20 bytes)
                        header = struct.unpack("<5I", f.read(20))
                        magic, version, n, d, dtype_code = header
                        if magic == 0x49534F43:  # "ISOC"
                            expected_bytes = n * d * 4
                            raw = f.read(expected_bytes)
                            if len(raw) == expected_bytes:
                                vectors = np.frombuffer(raw, dtype=np.float32).reshape(n, d).copy()
                            else:
                                logger.error(
                                    "Truncated vectors.bin for %s: expected %d bytes, got %d",
                                    safe_name, expected_bytes, len(raw),
                                )
                        else:
                            logger.error(
                                "Invalid vectors.bin magic for %s: 0x%08X",
                                safe_name, magic,
                            )
            except Exception as exc:
                logger.error("Failed to load vectors for %s: %s", safe_name, exc)

        if metadata_path.exists():
            try:
                meta_data = json.loads(metadata_path.read_text(encoding="utf-8"))
                metadata = meta_data.get("chunks", [])
            except Exception as exc:
                logger.error("Failed to load metadata for %s: %s", safe_name, exc)

        with self._lock:
            self._loaded_indices[safe_name] = {
                "vectors": vectors,
                "metadata": metadata,
            }

        logger.info(
            "Index loaded  name=%s  vectors=%d  metadata=%d",
            safe_name,
            vectors.shape[0] if vectors is not None else 0,
            len(metadata),
        )
        return True

    def unload_index(self, name: str) -> None:
        """Unload an index from memory."""
        safe_name = self._safe_index_name(name)
        with self._lock:
            self._loaded_indices.pop(safe_name, None)
        logger.info("Index unloaded  name=%s", safe_name)

    # ------------------------------------------------------------------
    # File Ingestion Pipeline
    # ------------------------------------------------------------------

    def ingest_files(
        self,
        index_name: str,
        file_paths: list[str],
        progress_callback: Optional[Callable] = None,
    ) -> IngestionStats:
        """
        Ingest files into an index.
        Pipeline: scan files -> extract text -> chunk -> embed -> store vectors.
        """
        stats = IngestionStats()
        t_start = time.perf_counter()

        safe_name = self._safe_index_name(index_name)

        # Validate index exists
        info = self._read_index_info(safe_name)
        if info is None:
            raise FileNotFoundError(f"Index '{index_name}' not found")

        # Get a DB connection for dedup checks and batch inserts
        conn = self._get_db()

        # Load existing data
        self.load_index(safe_name)
        with self._lock:
            idx_data = self._loaded_indices.get(safe_name)
        if idx_data is None:
            raise ValueError(f"Index '{safe_name}' failed to load")

        existing_vectors = idx_data["vectors"]
        existing_metadata = idx_data["metadata"]

        if existing_vectors is None:
            if np is None:
                raise RuntimeError(
                    "numpy is required for indexing operations. "
                    "Install with: pip install numpy"
                )
            existing_vectors = np.zeros((0, DEFAULT_VECTOR_DIM), dtype=np.float32)

        # Ensure model
        if not self.ensure_model():
            raise RuntimeError(
                "Could not load the AI embedding model. "
                "Please install onnxruntime and tokenizers: "
                "pip install onnxruntime tokenizers"
            )

        # Process each file
        texts_to_embed = []
        texts_meta = []
        doc_records = []  # For database tracking

        for file_path in file_paths:
            stats.files_scanned += 1
            path = Path(file_path)

            if not path.exists():
                logger.warning("File not found: %s", file_path)
                stats.files_failed += 1
                continue

            # Check file size limit
            try:
                file_size = path.stat().st_size
                if file_size > MAX_FILE_SIZE:
                    logger.warning("File too large: %s (%.1f MB)", path.name, file_size / (1024*1024))
                    stats.files_failed += 1
                    continue
                if file_size == 0:
                    logger.warning("Empty file: %s", path.name)
                    stats.files_failed += 1
                    continue
            except OSError:
                stats.files_failed += 1
                continue

            if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
                logger.warning("Unsupported format: %s", path.suffix)
                stats.files_failed += 1
                continue

            stats.files_accepted += 1

            try:
                # Extract text
                raw_text = self._extract_text(path)

                # Run pre_ingest plugin hook
                if self._plugin_manager:
                    try:
                        modified_text = self._plugin_manager.run_hook(
                            "pre_ingest", text=raw_text, file_path=str(path),
                            default=raw_text,
                        )
                        if modified_text and isinstance(modified_text, str):
                            raw_text = modified_text
                    except Exception:
                        pass

                if not raw_text or len(raw_text.strip()) < 3:
                    logger.warning("No text extracted from: %s", path.name)
                    stats.files_failed += 1
                    continue

                # Compute file hash for dedup
                file_hash = self._compute_file_hash(path)

                # Check for duplicate files (same hash already indexed)
                if file_hash:
                    try:
                        existing = conn.execute(
                            "SELECT doc_id FROM documents WHERE index_name = ? AND file_hash = ?",
                            (safe_name, file_hash),
                        ).fetchone()
                        if existing:
                            logger.info("Skipping duplicate file: %s (hash=%s)", path.name, file_hash[:12])
                            stats.files_scanned += 0  # already counted
                            stats.files_failed += 0
                            continue
                    except Exception:
                        pass

                # For code files, also generate a code structure summary
                code_summary = ""
                if path.suffix.lower() in {".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp", ".h", ".hpp", ".go", ".rs", ".rb"}:
                    try:
                        from desktop_app.code_parser import parse_code_file, generate_code_summary
                        parsed = parse_code_file(path)
                        if parsed and parsed.symbols:
                            code_summary = generate_code_summary(parsed)
                            logger.info("Parsed %d symbols from %s", len(parsed.symbols), path.name)
                    except Exception as exc:
                        logger.debug("Code parsing failed for %s: %s", path.name, exc)

                # Chunk
                chunks = self._chunk_text(raw_text, path.name)
                if not chunks:
                    stats.files_failed += 1
                    continue

                # If we have a code summary, prepend it as a special chunk
                if code_summary:
                    chunks.insert(0, {
                        "text": code_summary,
                        "label": f"{path.name} (code structure)",
                        "page_number": 0,
                    })

                stats.files_processed += 1
                stats.total_chunks += len(chunks)

                start_meta_idx = len(texts_meta)
                for chunk in chunks:
                    texts_to_embed.append(chunk["text"])
                    texts_meta.append({
                        "full_text": chunk["text"],
                        "text_preview": chunk["text"][:200],
                        "source_file": path.name,
                        "source_path": str(path),
                        "source_label": chunk.get("label", ""),
                        "format_category": path.suffix.lower().lstrip("."),
                        "word_count": len(chunk["text"].split()),
                        "page_number": chunk.get("page_number", 0),
                    })

                # Track document for database
                doc_records.append({
                    "doc_id": str(uuid.uuid4()),
                    "index_name": safe_name,
                    "file_path": str(path),
                    "file_hash": file_hash,
                    "format_category": path.suffix.lower().lstrip("."),
                    "chunk_count": len(chunks),
                    "word_count": len(raw_text.split()),
                    "meta_start_idx": start_meta_idx,
                    "meta_end_idx": len(texts_meta),
                })

                if progress_callback:
                    progress_callback(
                        stats.files_processed,
                        len(file_paths),
                        path.name,
                    )

            except Exception as exc:
                logger.error("Failed to process %s: %s", path.name, exc)
                stats.files_failed += 1

        # Batch embed
        if texts_to_embed:
            assert np is not None  # guaranteed by embed_batch succeeding
            vectors = self.embed_batch(texts_to_embed)
            if vectors is None:
                raise RuntimeError(
                    "Embedding failed — could not generate vectors for "
                    f"{len(texts_to_embed)} text chunks. Check that onnxruntime "
                    "is installed and the embedding model is valid."
                )
            if existing_vectors.shape[0] == 0:
                new_vectors = vectors
            else:
                new_vectors = np.vstack([existing_vectors, vectors])

            new_metadata = existing_metadata + texts_meta

            # Save to disk
            self._save_index_data(safe_name, new_vectors, new_metadata)

            # Update memory — hold lock to prevent concurrent write-back
            with self._lock:
                # Re-read idx_data in case another thread already updated it
                current = self._loaded_indices.get(safe_name)
                if current is not None:
                    cur_vec = current["vectors"]
                    cur_meta = current["metadata"]
                    if cur_vec is not None and cur_vec.shape[0] != existing_vectors.shape[0]:
                        # Another ingestion happened concurrently — merge
                        new_vectors = np.vstack([cur_vec, vectors]) if cur_vec.shape[0] > 0 else vectors
                        new_metadata = cur_meta + texts_meta
                    current["vectors"] = new_vectors
                    current["metadata"] = new_metadata

            # Update info
            info["vector_count"] = len(new_metadata)
            info["updated_at"] = datetime.now(timezone.utc).isoformat()
            self._write_index_info(safe_name, info)

            stats.total_vectors = vectors.shape[0]

        # Persist document records to database (batch insert)
        if doc_records:
            conn = self._get_db()
            try:
                for doc in doc_records:
                    conn.execute(
                        """INSERT OR REPLACE INTO documents
                           (doc_id, index_name, file_path, file_hash, format_category,
                            chunk_count, word_count, status)
                           VALUES (?, ?, ?, ?, ?, ?, ?, 'indexed')""",
                        (doc["doc_id"], doc["index_name"], doc["file_path"],
                         doc["file_hash"], doc["format_category"],
                         doc["chunk_count"], doc["word_count"]),
                    )
                conn.commit()
            except Exception:
                pass

        stats.elapsed_seconds = time.perf_counter() - t_start
        self._record_event("document_ingested", metadata={
            "index_name": safe_name,
            "files_processed": stats.files_processed,
            "chunks_created": stats.total_chunks,
        })

        logger.info(
            "Ingestion complete  index=%s  files=%d  chunks=%d  elapsed=%.1fs",
            safe_name,
            stats.files_processed,
            stats.total_chunks,
            stats.elapsed_seconds,
        )
        return stats

    def _compute_file_hash(self, file_path: Path) -> str:
        """Compute SHA-256 hash of a file for deduplication."""
        sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for block in iter(lambda: f.read(65536), b""):
                    sha256.update(block)
            return sha256.hexdigest()
        except Exception:
            return ""

    def _extract_text(self, file_path: Path) -> str:
        """Extract text from a file based on its format."""
        ext = file_path.suffix.lower()

        # Plain text
        if ext in {".txt", ".md", ".log", ".rst", ".cfg", ".ini", ".py", ".cpp",
                    ".c", ".h", ".js", ".ts", ".java", ".go", ".rs", ".rb",
                    ".csv", ".json", ".tsv", ".html", ".htm"}:
            for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
                try:
                    return file_path.read_text(encoding=encoding)
                except (UnicodeDecodeError, LookupError):
                    continue
            return file_path.read_text(encoding="utf-8", errors="replace")

        # PDF — multi-strategy extraction for better quality
        if ext == ".pdf":
            return self._extract_pdf(file_path)

        # DOCX
        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(file_path))
                lines = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n".join(lines)
            except Exception as exc:
                logger.error("DOCX extraction failed for %s: %s", file_path.name, exc)
                return ""

        # PPTX
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                slides = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():  # type: ignore[union-attr]
                            slide_text.append(shape.text)  # type: ignore[union-attr]
                    if slide_text:
                        slides.append("\n".join(slide_text))
                return "\n\n---\n\n".join(slides)
            except Exception as exc:
                logger.error("PPTX extraction failed for %s: %s", file_path.name, exc)
                return ""

        # XLSX / XLS
        if ext in {".xlsx", ".xls"}:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
                sheets = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                        if cells:
                            rows.append(" | ".join(cells))
                    if rows:
                        sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
                wb.close()
                return "\n\n".join(sheets)
            except Exception as exc:
                logger.error("XLSX extraction failed for %s: %s", file_path.name, exc)
                return ""

        # Email
        if ext == ".eml":
            try:
                import email
                with open(file_path, "rb") as f:
                    msg = email.message_from_bytes(f.read())
                parts = []
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            parts.append(payload.decode("utf-8", errors="replace"))  # type: ignore[union-attr]
                return "\n".join(parts)
            except Exception as exc:
                logger.error("EML extraction failed for %s: %s", file_path.name, exc)
                return ""

        # Fallback: try reading as text
        try:
            text = file_path.read_text(encoding="utf-8", errors="replace")
            if text and len(text.strip()) > 10:
                return text
        except Exception:
            pass

        # OCR fallback for images
        if ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"}:
            try:
                from desktop_app.ocr import ocr_image
                return ocr_image(file_path)
            except Exception as exc:
                logger.debug("OCR fallback failed for %s: %s", file_path.name, exc)

        return ""

    def extract_text(self, file_path) -> str:
        """Public wrapper around _extract_text().

        Args:
            file_path: Path-like pointing to a supported document.

        Returns:
            Extracted text content, or empty string on failure.
        """
        return self._extract_text(Path(file_path) if not isinstance(file_path, Path) else file_path)

    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF using multi-strategy approach.
        
        Strategy 1: PyMuPDF (fitz) — fast, good for simple PDFs
        Strategy 2: pdfplumber — better for tables and complex layouts
        
        Falls back if text yield is suspiciously low.
        """
        text_result = ""
        import fitz
        
        # Check if encrypted
        try:
            doc = fitz.open(str(file_path))
        except Exception as exc:
            logger.error("PDF open failed for %s: %s", file_path.name, exc)
            return ""
        
        if doc.is_encrypted:
            doc.close()
            return ""
        
        total_pages = len(doc)
        file_size = file_path.stat().st_size
        
        # --- Strategy 1: PyMuPDF ---
        try:
            pages = []
            for i, page in enumerate(doc):
                text = page.get_text("text")
                if text and text.strip():
                    pages.append(f"<<<PAGE:{i + 1}>>>{text}")
            doc.close()
            
            text_result = "\n\n".join(pages)
            
            # Heuristic: if file is > 50KB but we got < 500 chars per page,
            # the PDF might have complex layout (tables, columns, scanned images)
            avg_chars_per_page = len(text_result) / max(total_pages, 1)
            expected_chars = file_size / 50  # rough estimate: 50 bytes per char of text
            
            if len(text_result.strip()) > 200 or avg_chars_per_page > 100:
                return text_result
            
            # Low yield — try pdfplumber
            logger.info(
                "PyMuPDF low yield for %s (%.0f chars/page), trying pdfplumber",
                file_path.name, avg_chars_per_page,
            )
        except Exception as exc:
            logger.warning("PyMuPDF extraction failed for %s: %s", file_path.name, exc)
        
        # --- Strategy 2: pdfplumber ---
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber not installed, using PyMuPDF result")
            return text_result
        
        try:
            pages = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    # Extract text with layout preservation
                    text = page.extract_text(layout=True) or ""
                    
                    # Also extract tables and append them
                    tables = page.extract_tables()
                    if tables:
                        for table in tables:
                            # Format table as text
                            table_text = self._format_table(table)
                            if table_text:
                                text += "\n\n[TABLE]\n" + table_text
                    
                    if text and text.strip():
                        pages.append(f"<<<PAGE:{i + 1}>>>{text}")
            
            result = "\n\n".join(pages)
            if result.strip():
                logger.info("pdfplumber extracted %d chars from %s", len(result), file_path.name)
                return result
            else:
                # pdfplumber also failed — return whatever PyMuPDF got
                return text_result
        except Exception as exc:
            logger.error("pdfplumber extraction failed for %s: %s", file_path.name, exc)

        # --- Strategy 3: OCR fallback for scanned PDFs ---
        logger.info("Text extraction low yield for %s, trying OCR", file_path.name)
        try:
            from desktop_app.ocr import ocr_pdf
            ocr_text = ocr_pdf(file_path)
            if ocr_text and len(ocr_text.strip()) > 50:
                return ocr_text
        except Exception as exc:
            logger.debug("OCR fallback failed for %s: %s", file_path.name, exc)

        return text_result

    @staticmethod
    def _format_table(table: list[list]) -> str:
        """Format a pdfplumber table as readable text."""
        if not table:
            return ""
        
        # Clean cells
        rows = []
        for row in table:
            cells = [str(c).strip() if c else "" for c in row]
            if any(cells):  # Skip empty rows
                rows.append(" | ".join(cells))
        
        return "\n".join(rows)

    def _chunk_text(self, text: str, source_name: str) -> list[dict]:
        """Chunk text into segments for embedding. Tracks PDF page numbers via <<<PAGE:N>>> markers."""
        import re

        if not text or not text.strip():
            return []

        # Clean page markers from text but track current page number
        page_markers = re.findall(r'<<<PAGE:(\d+)>>>', text)
        text_clean = re.sub(r'<<<PAGE:\d+>>>', '', text)

        # Build a mapping: for each word index, what page is it on?
        # Split into segments by page markers to track boundaries
        segments = []
        if page_markers:
            # Split text at page boundaries
            parts = re.split(r'<<<PAGE:\d+>>>', text)
            pages = [0] + [int(p) for p in page_markers]
            for i, part in enumerate(parts):
                if part.strip():
                    segments.append((pages[i], part.strip()))
        else:
            segments = [(0, text_clean.strip())]

        if not segments:
            return []

        # Now chunk across all segments, tracking which page each word belongs to
        chunks = []
        current_words = []       # list of (word, page_number)
        current_count = 0
        chunk_idx = 0

        def _flush():
            nonlocal current_words, current_count, chunk_idx
            if not current_words:
                return
            words_only = [w for w, _ in current_words]
            # Majority page number for this chunk
            page_counts: dict[int, int] = {}
            for _, pn in current_words:
                page_counts[pn] = page_counts.get(pn, 0) + 1
            dominant_page = max(page_counts, key=lambda k: page_counts[k]) if page_counts else 0

            chunk_text = " ".join(words_only)
            chunks.append({
                "text": chunk_text,
                "label": f"{source_name} (chunk {chunk_idx})",
                "page_number": dominant_page,
            })
            chunk_idx += 1

            # Overlap
            if len(current_words) > DEFAULT_OVERLAP:
                current_words = current_words[-DEFAULT_OVERLAP:]
                current_count = len(current_words)
            else:
                current_words = []
                current_count = 0

        for page_num, segment_text in segments:
            # Split into sentences within this segment
            sentences = re.split(r'(?<=[.!?])\s+|(?<=\n)\s*\n\s*(?=\S)', segment_text)
            sentences = [s.strip() for s in sentences if s.strip()]
            if not sentences:
                sentences = [segment_text]

            for sentence in sentences:
                words = sentence.split()
                if not words:
                    continue

                if current_count > 0 and current_count + len(words) > DEFAULT_CHUNK_SIZE:
                    _flush()

                for w in words:
                    current_words.append((w, page_num))
                current_count += len(words)

        _flush()
        return chunks

    def _save_index_data(self, name: str, vectors: np.ndarray, metadata: list) -> None:  # type: ignore[name-defined]
        """Save vectors and metadata to disk.

        Vector binary format:
          magic   (4 bytes, uint32) = 0x49534F43 ("ISOC")
          version (4 bytes, uint32) = 1
          count   (4 bytes, uint32) = number of vectors
          dim     (4 bytes, uint32) = vector dimension
          dtype   (4 bytes, uint32) = 0 for float32
          data    (count * dim * 4 bytes) = raw float32 vectors
        """
        import struct
        import tempfile

        idx_dir = self._indices_dir / name
        idx_dir.mkdir(parents=True, exist_ok=True)

        # Save vectors in custom binary format using atomic write
        vectors_path = idx_dir / "vectors.bin"
        n, d = vectors.shape

        header = struct.pack(
            "<5I",
            0x49534F43,   # magic: "ISOC" in little-endian
            1,            # version
            n,            # count
            d,            # dimension
            0,            # dtype: float32
        )

        # Write to temp file first, then atomically rename
        tmp_vectors = vectors_path.with_suffix(".bin.tmp")
        try:
            with open(str(tmp_vectors), "wb") as f:
                f.write(header)
                vectors.tofile(f)
            tmp_vectors.replace(vectors_path)
        except Exception:
            tmp_vectors.unlink(missing_ok=True)
            raise

        # Save metadata atomically
        metadata_path = idx_dir / "metadata.json"
        tmp_meta = metadata_path.with_suffix(".json.tmp")
        try:
            tmp_meta.write_text(
                json.dumps({"chunks": metadata}, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            tmp_meta.replace(metadata_path)
        except Exception:
            tmp_meta.unlink(missing_ok=True)
            raise

        logger.info(
            "Saved index data  name=%s  vectors=%d  dim=%d  size=%.1fMB",
            name, n, d,
            vectors_path.stat().st_size / (1024 * 1024),
        )

    # ------------------------------------------------------------------
    # Semantic Search
    # ------------------------------------------------------------------

    def search(
        self,
        index_name: str,
        query: str,
        top_k: int = 10,
    ) -> list[SearchResult]:
        """
        Perform semantic search across an index.

        Uses cosine similarity between the query embedding and all
        stored vectors. Returns the top-k results ranked by score.

        Args:
            index_name: Name of the index to search.
            query: Natural language query string.
            top_k: Maximum number of results to return.

        Returns:
            List of SearchResult objects, sorted by score descending.

        Raises:
            ValueError: If the index doesn't exist, has no vectors,
                        or the query is too short.
        """
        if not query or len(query.strip()) < 2:
            raise ValueError("Query must be at least 2 characters")

        top_k = max(1, min(top_k, 100))
        safe_name = self._safe_index_name(index_name)

        # Ensure index is loaded
        self.load_index(safe_name)
        with self._lock:
            idx_data = self._loaded_indices.get(safe_name)
        if idx_data is None:
            raise ValueError(f"Index '{index_name}' not found")

        vectors = idx_data["vectors"]
        metadata = idx_data["metadata"]

        if vectors is None or vectors.shape[0] == 0:
            raise ValueError(f"Index '{safe_name}' has no indexed vectors yet. "
                             "Upload and index files first.")

        # Embed the query
        query_vec = self.embed_text(query)
        if query_vec is None:
            raise ValueError("Failed to embed query. Check that the embedding model is available.")

        if np is None:
            raise ValueError("numpy is required for search operations")

        # Compute cosine similarity
        # Since embeddings are normalized, dot product = cosine similarity
        scores = np.dot(vectors, query_vec)

        # Get top-k indices
        top_indices = np.argsort(scores)[::-1][:top_k]

        results = []
        query_lower = query.lower().split()
        for rank, idx in enumerate(top_indices):
            score = float(scores[idx])

            # Skip near-zero results
            if score < 0.05:
                continue

            meta = metadata[idx] if idx < len(metadata) else {}
            full_text = meta.get("full_text", meta.get("text_preview", ""))
            page_number = meta.get("page_number", 0)

            # Find the best matching segment in full text around query words
            matched_segment = self._find_context_snippet(full_text, query_lower, window=80)

            results.append(SearchResult(
                rank=rank + 1,
                text=meta.get("text_preview", ""),
                source_file=meta.get("source_file", "unknown"),
                source_label=meta.get("source_label", ""),
                score=score,
                chunk_index=int(idx),
                format_category=meta.get("format_category", ""),
                page_number=page_number,
                full_text=full_text,
                matched_segment=matched_segment,
            ))

        # Record analytics
        self._record_event("search", metadata={
            "index_name": safe_name,
            "query_length": len(query),
            "results_count": len(results),
            "top_score": round(float(results[0].score), 4) if results else 0,
        })

        logger.info(
            "Search  index=%s  query='%s'  results=%d  top_score=%.4f",
            safe_name, query[:50], len(results),
            results[0].score if results else 0,
        )
        return results

    # ------------------------------------------------------------------
    # Hybrid Search (Semantic + Keyword)
    # ------------------------------------------------------------------

    def _keyword_score(self, text: str, query: str) -> float:
        """Compute a simple BM25-style keyword score for text vs query.

        Uses TF-IDF-like scoring: term frequency * inverse document frequency.
        No external dependencies — pure Python.
        """
        if not text or not query:
            return 0.0

        query_terms = query.lower().split()
        text_lower = text.lower()
        text_terms = text_lower.split()

        if not query_terms or not text_terms:
            return 0.0

        total_terms = len(text_terms)
        scores = 0.0

        for qt in query_terms:
            if len(qt) < 2:
                continue

            # Term frequency in document
            tf = text_terms.count(qt)
            if tf == 0:
                continue

            # Normalize TF
            tf_norm = tf / total_terms

            # Simple IDF approximation (log of inverse doc freq)
            # We use a fixed IDF of log(N/df) where N=1000 (typical corpus), df=10
            idf = math.log(1000 / 10)

            scores += tf_norm * idf

            # Bonus for exact phrase match
            if qt in text_lower:
                scores += 0.1

        # Normalize by number of query terms that matched
        matched = sum(1 for qt in query_terms if len(qt) >= 2 and qt in text_terms)
        if matched > 0:
            scores = scores * (matched / len([qt for qt in query_terms if len(qt) >= 2]))

        return min(scores, 1.0)

    def _hybrid_rank(
        self,
        results: list,
        query: str,
        semantic_weight: float = 0.7,
        keyword_weight: float = 0.3,
    ) -> list:
        """Re-rank search results by blending semantic and keyword scores.

        Args:
            results: List of SearchResult objects.
            query: Original query string.
            semantic_weight: Weight for semantic (embedding) score (0-1).
            keyword_weight: Weight for keyword (BM25) score (0-1).

        Returns:
            Re-ranked list of SearchResult with updated scores.
        """
        if not results:
            return results

        # Compute keyword scores for each result
        for r in results:
            # Score against both the full text and the text preview
            kw_full = self._keyword_score(r.full_text or "", query)
            kw_preview = self._keyword_score(r.text or "", query)
            r._keyword_score = max(kw_full, kw_preview)

            # Blended score
            r._blended_score = (
                semantic_weight * r.score +
                keyword_weight * r._keyword_score
            )

        # Re-sort by blended score
        results.sort(key=lambda r: r._blended_score, reverse=True)

        # Update rank numbers
        for i, r in enumerate(results):
            r.rank = i + 1

        return results

    # ------------------------------------------------------------------
    # Context Snippet
    # ------------------------------------------------------------------

    @staticmethod
    def _find_context_snippet(text: str, query_words: list[str], window: int = 80) -> str:
        """Find the best ~window-word segment of *text* that contains the most query words.
        Returns the snippet with \x00 markers around each matched query word.
        """
        if not text or not query_words:
            return ""
        words = text.split()
        if len(words) <= window * 2:
            # Text is short enough — highlight in place
            return _highlight_words(" ".join(words), query_words)

        # Score each window position by how many query words it contains
        best_pos = 0
        best_count = 0
        text_lower = [w.lower() for w in words]
        query_set = set(query_words)

        step = max(window // 4, 10)
        for start in range(0, len(words) - window + 1, step):
            segment = text_lower[start:start + window]
            hits = sum(1 for w in segment if w in query_set)
            if hits > best_count:
                best_count = hits
                best_pos = start

        if best_count == 0:
            best_pos = 0  # just show the beginning

        end = min(best_pos + window, len(words))
        snippet_words = words[best_pos:end]
        prefix = "\u2026 " if best_pos > 0 else ""
        suffix = " \u2026" if end < len(words) else ""
        return prefix + _highlight_words(" ".join(snippet_words), query_words) + suffix

    # ------------------------------------------------------------------
    # RAG: Search All Indexes
    # ------------------------------------------------------------------

    def search_all_indexes(
        self,
        query: str,
        top_k: int = 5,
    ) -> list[SearchResult]:
        """
        Search across ALL loaded indexes and return the top-k best
        results ranked by score. Used by the RAG pipeline to build
        context for the LLM.

        Args:
            query: Natural language query.
            top_k: Total results to return (across all indexes).

        Returns:
            List of SearchResult, sorted by score descending.
        """
        if not query or len(query.strip()) < 2:
            return []

        all_results: list[SearchResult] = []

        # Get list of all index directories
        index_names = self._list_index_dirs()
        if not index_names:
            return []

        # Ensure embedding model is ready
        self.ensure_model()

        # Run pre_search plugin hook
        if self._plugin_manager:
            try:
                modified_query = self._plugin_manager.run_hook(
                    "pre_search", query=query, default=query,
                )
                if modified_query and isinstance(modified_query, str):
                    query = modified_query
            except Exception:
                pass

        # Embed query ONCE for all indexes (performance fix)
        query_vec = self.embed_text(query)
        if query_vec is None:
            logger.warning("Failed to embed query for search_all_indexes")
            return []

        if np is None:
            return []

        query_lower = query.lower().split()

        for name in index_names:
            try:
                safe_name = self._safe_index_name(name)
                self.load_index(safe_name)
                with self._lock:
                    idx_data = self._loaded_indices.get(safe_name)
                if idx_data is None:
                    continue

                vectors = idx_data["vectors"]
                metadata = idx_data["metadata"]

                if vectors is None or vectors.shape[0] == 0:
                    continue

                # Compute cosine similarity with pre-computed query vector
                scores = np.dot(vectors, query_vec)
                top_indices = np.argsort(scores)[::-1][:10]

                for rank, idx in enumerate(top_indices):
                    score = float(scores[idx])
                    if score < 0.05:
                        continue

                    meta = metadata[idx] if idx < len(metadata) else {}
                    full_text = meta.get("full_text", meta.get("text_preview", ""))
                    page_number = meta.get("page_number", 0)
                    matched_segment = self._find_context_snippet(full_text, query_lower, window=80)

                    all_results.append(SearchResult(
                        rank=len(all_results) + 1,
                        text=meta.get("text_preview", ""),
                        source_file=meta.get("source_file", "unknown"),
                        source_label=meta.get("source_label", ""),
                        score=score,
                        chunk_index=int(idx),
                        format_category=meta.get("format_category", ""),
                        page_number=page_number,
                        full_text=full_text,
                        matched_segment=matched_segment,
                    ))

            except Exception as exc:
                logger.debug("Index '%s' search skipped: %s", name, exc)
                continue

        # Sort by semantic score first, take more than top_k for hybrid re-ranking
        all_results.sort(key=lambda r: r.score, reverse=True)

        # Apply hybrid re-ranking (blend semantic + keyword scores)
        if len(all_results) > 0:
            all_results = self._hybrid_rank(all_results[:top_k * 3], query)

        return all_results[:top_k]

    def _build_stats_summary(self) -> str:
        """Build a comprehensive data summary for the AI.

        Includes:
        - Index-level stats (files, chunks, vectors, HNSW params)
        - Complete file listing per index (name, format, word count, chunk count)
        - Watched folders
        - Total library statistics

        This is always injected into the RAG context so the AI can answer
        any question about the user's data.
        """
        try:
            indexes = self.list_indexes()
            if not indexes:
                return ""

            conn = self._get_db()
            total_files = 0
            total_chunks = 0
            total_words = 0

            parts = ["[IsoCortex Library Data]"]

            for idx in indexes:
                # Per-index aggregate stats
                row = conn.execute(
                    "SELECT COUNT(*), SUM(chunk_count), SUM(word_count) "
                    "FROM documents WHERE index_name = ? AND status = 'indexed'",
                    (idx.name,),
                ).fetchone()
                f_count = row[0] or 0
                c_count = row[1] or 0
                w_count = row[2] or 0
                total_files += f_count
                total_chunks += c_count
                total_words += w_count

                parts.append(
                    f"\nIndex: '{idx.name}' "
                    f"(created {idx.created_at[:10] if idx.created_at else 'unknown'})"
                )
                if idx.description:
                    parts.append(f"  Description: {idx.description}")
                parts.append(
                    f"  Stats: {f_count} files, {c_count} chunks, "
                    f"{idx.vector_count} vectors in HNSW graph, "
                    f"{w_count:,} words total"
                )

                # List every file in this index
                if f_count > 0:
                    file_rows = conn.execute(
                        "SELECT file_path, format_category, chunk_count, word_count "
                        "FROM documents WHERE index_name = ? AND status = 'indexed' "
                        "ORDER BY created_at",
                        (idx.name,),
                    ).fetchall()

                    parts.append("  Files:")
                    for fr in file_rows:
                        fpath = fr[0]
                        # Show just the filename for readability
                        fname = os.path.basename(fpath) if fpath else "unknown"
                        fmt = fr[1] or "unknown"
                        chunks = fr[2] or 0
                        words = fr[3] or 0
                        parts.append(
                            f"    - {fname} [{fmt}] — {chunks} chunks, {words:,} words"
                        )

            # Watched folders
            watch_rows = conn.execute(
                "SELECT folder_path, index_name, is_active FROM watched_folders"
            ).fetchall()
            if watch_rows:
                parts.append("\nWatched Folders:")
                for wr in watch_rows:
                    status = "active" if wr[2] else "inactive"
                    parts.append(f"  - {wr[0]} -> index '{wr[1]}' ({status})")

            # Totals
            parts.append(
                f"\nTotals: {total_files} files, {total_chunks} chunks, "
                f"{total_words:,} words across {len(indexes)} index(es)"
            )

            return "\n".join(parts)
        except Exception as exc:
            logger.debug("Could not build stats summary: %s", exc)
            return ""

    def build_rag_context(self, query: str, top_k: int = 5) -> tuple[str, list[dict]]:
        """
        Retrieve relevant chunks from all indexes and format them
        as a RAG context string for the LLM.

        Always includes a brief stats summary so the AI can answer
        metadata questions (how many files indexed, etc.).
        """
        results = self.search_all_indexes(query, top_k=top_k)

        # ── Build stats summary (always included) ────────────────
        stats_summary = self._build_stats_summary()

        if not results and not stats_summary:
            return "", []

        context_parts = []
        sources = []
        seen_texts = set()

        for i, r in enumerate(results):
            # Deduplicate by text content
            text_key = r.text[:100] if r.text else ""
            if text_key in seen_texts:
                continue
            seen_texts.add(text_key)

            chunk_text = r.full_text or r.text or r.matched_segment or ""
            if not chunk_text:
                continue

            source_file = r.source_file or "unknown"
            # Use just the filename, not full path
            if "/" in source_file or "\\" in source_file:
                source_file = os.path.basename(source_file)

            # Get keyword score if available from hybrid ranking
            kw_score = getattr(r, '_keyword_score', 0) or 0
            blended_score = getattr(r, '_blended_score', r.score) or r.score

            source_info = {
                "index": i + 1,
                "file": source_file,
                "score": round(blended_score, 3),
                "keyword_score": round(kw_score, 3),
                "page": r.page_number or 0,
                "format": r.format_category or "",
                "full_text": chunk_text,
                "chunk_text": chunk_text,
            }
            sources.append(source_info)

            context_parts.append(
                f"[Source {i + 1}] ({source_file}"
                + (f", page {r.page_number}" if r.page_number else "")
                + f"):\n{chunk_text}"
            )

        # Budget: approximately 2000 tokens (~8000 chars) for document chunks
        context_str = "\n\n---\n\n".join(context_parts)
        if len(context_str) > 8000:
            # Truncate from the end, keeping whole sources
            truncated = []
            length = 0
            for part in context_parts:
                if length + len(part) > 8000:
                    break
                truncated.append(part)
                length += len(part)
            context_str = "\n\n---\n\n".join(truncated)

        # Prepend stats summary so the AI always has metadata context
        if stats_summary:
            context_str = stats_summary + "\n\n" + context_str

        return context_str, sources

    # ------------------------------------------------------------------
    # Chat Persistence
    # ------------------------------------------------------------------

    def create_conversation(self) -> str:
        """Create a new chat conversation and return its ID."""
        conv_id = str(uuid.uuid4())
        user_id = self._current_user["user_id"] if self._current_user else None
        conn = self._get_db()
        conn.execute(
            "INSERT INTO chat_conversations (conversation_id, user_id, title) VALUES (?, ?, ?)",
            (conv_id, user_id, "New Chat"),
        )
        conn.commit()
        logger.info("Created conversation %s", conv_id)
        return conv_id

    def save_message(
        self,
        conversation_id: str,
        role: str,
        content: str,
        sources: Optional[list[dict]] = None,
        token_count: int = 0,
    ) -> str:
        """Save a chat message and return its ID."""
        msg_id = str(uuid.uuid4())
        sources_json = json.dumps(sources or [], default=str)
        conn = self._get_db()
        conn.execute(
            "INSERT INTO chat_messages (message_id, conversation_id, role, content, sources_json, token_count) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (msg_id, conversation_id, role, content, sources_json, token_count),
        )
        # Update conversation timestamp
        conn.execute(
            "UPDATE chat_conversations SET updated_at = datetime('now') WHERE conversation_id = ?",
            (conversation_id,),
        )
        conn.commit()
        return msg_id

    def get_conversation_messages(self, conversation_id: str) -> list[dict]:
        """Get all messages for a conversation, ordered by creation time."""
        conn = self._get_db()
        rows = conn.execute(
            "SELECT message_id, role, content, sources_json, token_count, created_at "
            "FROM chat_messages WHERE conversation_id = ? ORDER BY created_at ASC",
            (conversation_id,),
        ).fetchall()
        messages = []
        for row in rows:
            msg_id, role, content, sources_json, tok_count, created_at = row
            sources = json.loads(sources_json) if sources_json else []
            messages.append({
                "message_id": msg_id,
                "role": role,
                "content": content,
                "sources": sources,
                "token_count": tok_count,
                "created_at": created_at,
            })
        return messages

    def list_conversations(self, limit: int = 50) -> list[dict]:
        """List recent conversations, newest first."""
        conn = self._get_db()
        rows = conn.execute(
            "SELECT conversation_id, title, created_at, updated_at "
            "FROM chat_conversations ORDER BY updated_at DESC LIMIT ?",
            (limit,),
        ).fetchall()
        return [
            {
                "conversation_id": row[0],
                "title": row[1],
                "created_at": row[2],
                "updated_at": row[3],
            }
            for row in rows
        ]

    def update_conversation_title(self, conversation_id: str, title: str) -> None:
        """Update the title of a conversation."""
        conn = self._get_db()
        conn.execute(
            "UPDATE chat_conversations SET title = ? WHERE conversation_id = ?",
            (title[:100], conversation_id),
        )
        conn.commit()

    def delete_conversation(self, conversation_id: str) -> None:
        """Delete a conversation and all its messages."""
        conn = self._get_db()
        conn.execute(
            "DELETE FROM chat_messages WHERE conversation_id = ?",
            (conversation_id,),
        )
        conn.execute(
            "DELETE FROM chat_conversations WHERE conversation_id = ?",
            (conversation_id,),
        )
        conn.commit()
        logger.info("Deleted conversation %s", conversation_id)

    # ------------------------------------------------------------------
    # Watch Folders
    # ------------------------------------------------------------------

    def add_watched_folder(self, folder_path: str, index_name: str = "default") -> str:
        """Add a folder to the watch list. Returns folder_id."""
        folder_id = str(uuid.uuid4())
        conn = self._get_db()
        try:
            conn.execute(
                "INSERT INTO watched_folders (folder_id, folder_path, index_name) "
                "VALUES (?, ?, ?)",
                (folder_id, folder_path, index_name),
            )
            conn.commit()
        except Exception:
            logger.exception("Failed to add watched folder")
            raise
        logger.info("Added watched folder %s -> %s (%s)", folder_id, folder_path, index_name)
        return folder_id

    def remove_watched_folder(self, folder_id: str) -> bool:
        """Remove a watched folder by ID. Returns True if found and removed."""
        conn = self._get_db()
        cursor = conn.execute(
            "DELETE FROM watched_folders WHERE folder_id = ?",
            (folder_id,),
        )
        conn.commit()
        removed = cursor.rowcount > 0
        if removed:
            logger.info("Removed watched folder %s", folder_id)
        return removed

    def get_watched_folders(self) -> list[dict]:
        """Return all watched folders."""
        conn = self._get_db()
        rows = conn.execute(
            "SELECT folder_id, folder_path, index_name, is_active, created_at "
            "FROM watched_folders ORDER BY created_at"
        ).fetchall()
        return [
            {
                "folder_id": r[0],
                "folder_path": r[1],
                "index_name": r[2],
                "is_active": bool(r[3]),
                "created_at": r[4],
            }
            for r in rows
        ]

    def toggle_watched_folder(self, folder_id: str, active: bool) -> None:
        """Enable or disable a watched folder."""
        conn = self._get_db()
        conn.execute(
            "UPDATE watched_folders SET is_active = ? WHERE folder_id = ?",
            (int(active), folder_id),
        )
        conn.commit()
        logger.info("Toggled watched folder %s -> active=%s", folder_id, active)

    # ------------------------------------------------------------------
    # Search History
    # ------------------------------------------------------------------

    def save_search(self, query: str, results_count: int = 0) -> str:
        """Save a search query to history. Returns history_id."""
        history_id = str(uuid.uuid4())
        conn = self._get_db()
        conn.execute(
            "INSERT INTO search_history (history_id, query, results_count) VALUES (?, ?, ?)",
            (history_id, query, results_count),
        )
        conn.commit()
        return history_id

    def get_search_history(self, limit: int = 50) -> list[dict]:
        """Return recent search history, newest first."""
        conn = self._get_db()
        rows = conn.execute(
            "SELECT history_id, query, results_count, created_at "
            "FROM search_history ORDER BY created_at DESC LIMIT ?",
            (limit,),
        ).fetchall()
        return [
            {
                "history_id": r[0],
                "query": r[1],
                "results_count": r[2],
                "created_at": r[3],
            }
            for r in rows
        ]

    def clear_search_history(self) -> int:
        """Clear all search history. Returns count of deleted rows."""
        conn = self._get_db()
        cursor = conn.execute("DELETE FROM search_history")
        conn.commit()
        count = cursor.rowcount
        logger.info("Cleared %d search history entries", count)
        return count

    # ------------------------------------------------------------------
    # Bookmarks
    # ------------------------------------------------------------------

    def add_bookmark(
        self,
        query: str,
        file_path: str = "",
        chunk_text: str = "",
        score: float = 0,
        note: str = "",
    ) -> str:
        """Bookmark a search result. Returns bookmark_id."""
        bookmark_id = str(uuid.uuid4())
        conn = self._get_db()
        conn.execute(
            "INSERT INTO bookmarks (bookmark_id, query, file_path, chunk_text, score, note) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (bookmark_id, query, file_path, chunk_text, score, note),
        )
        conn.commit()
        logger.info("Added bookmark %s", bookmark_id)
        return bookmark_id

    def remove_bookmark(self, bookmark_id: str) -> bool:
        """Remove a bookmark. Returns True if found."""
        conn = self._get_db()
        cursor = conn.execute(
            "DELETE FROM bookmarks WHERE bookmark_id = ?",
            (bookmark_id,),
        )
        conn.commit()
        removed = cursor.rowcount > 0
        if removed:
            logger.info("Removed bookmark %s", bookmark_id)
        return removed

    def get_bookmarks(self, limit: int = 100) -> list[dict]:
        """Return all bookmarks, newest first."""
        conn = self._get_db()
        rows = conn.execute(
            "SELECT bookmark_id, history_id, query, file_path, chunk_text, score, note, created_at "
            "FROM bookmarks ORDER BY created_at DESC LIMIT ?",
            (limit,),
        ).fetchall()
        return [
            {
                "bookmark_id": r[0],
                "history_id": r[1],
                "query": r[2],
                "file_path": r[3],
                "chunk_text": r[4],
                "score": r[5],
                "note": r[6],
                "created_at": r[7],
            }
            for r in rows
        ]

    def update_bookmark_note(self, bookmark_id: str, note: str) -> None:
        """Update the note on a bookmark."""
        conn = self._get_db()
        conn.execute(
            "UPDATE bookmarks SET note = ? WHERE bookmark_id = ?",
            (note, bookmark_id),
        )
        conn.commit()

    # ------------------------------------------------------------------
    # Settings
    # ------------------------------------------------------------------

    def get_settings(self) -> dict:
        """Return all current settings and configuration."""
        return {
            "data_dir": str(self._data_dir),
            "indices_dir": str(self._indices_dir),
            "model_name": DEFAULT_MODEL_NAME,
            "vector_dim": DEFAULT_VECTOR_DIM,
            "hnsw": {
                "M": DEFAULT_HNSW_M,
                "ef_construction": DEFAULT_HNSW_EF_CONSTRUCTION,
                "ef_search": DEFAULT_HNSW_EF_SEARCH,
                "metric": "cosine",
            },
            "chunking": {
                "chunk_size": DEFAULT_CHUNK_SIZE,
                "overlap": DEFAULT_OVERLAP,
            },
            "batch_size": DEFAULT_BATCH_SIZE,
        }

    # ------------------------------------------------------------------
    # Data Reset
    # ------------------------------------------------------------------

    def reset_data(self) -> bool:
        """
        Reset all application data: indexes, documents, users, analytics.
        Used for the "Danger Zone" reset in Settings.

        Returns True on success.
        """
        try:
            # Unload all indexes from memory (thread-safe)
            with self._lock:
                self._loaded_indices.clear()

            # Close and invalidate the cached DB connection before deleting
            if self._db_conn is not None:
                try:
                    self._db_conn.close()
                except Exception:
                    pass
                self._db_conn = None

            # Close all thread-local DB connections
            for tid, tconn in list(getattr(self, '_thread_db_connections', {}).items()):
                try:
                    tconn.close()
                except Exception:
                    pass
            self._thread_db_connections = {}

            # Delete index data
            if self._indices_dir.exists():
                shutil.rmtree(self._indices_dir, ignore_errors=True)
            self._indices_dir.mkdir(parents=True, exist_ok=True)

            # Delete the database file (it will be recreated on next init)
            if self._db_path.exists():
                self._db_path.unlink()

            # Delete JWT secret
            secret_path = self._data_dir / ".jwt_secret"
            if secret_path.exists():
                secret_path.unlink()

            # Clear current user (thread-safe)
            with self._lock:
                self._current_user = None

            # Unload the model to free memory
            self._embed_model = None
            self._ort_session = None
            self._tokenizer = None

            # Reinitialize database (outside the lock to avoid deadlock)
            self._init_database()

            logger.info("All data has been reset successfully")
            return True
        except Exception as exc:
            logger.error("Failed to reset data: %s", exc)
            return False
